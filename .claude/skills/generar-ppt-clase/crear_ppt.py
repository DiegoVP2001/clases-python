#!/usr/bin/env python3
"""
crear_ppt.py — Generador del PPT de clase desde el spec enriquecido. (v5)

Mejoras v5:
- El PPT ahora termina en "Errores comunes". Guiada, Independiente, Ticket
  y Cierre se eliminaron del PPT — esas secciones se trabajan desde el Colab.
- Nueva capa de planificación (planificar_slides.py) que decide cuántos slides
  genera la ICN y con qué composición, usando presupuesto de densidad (filas).
- Slides ICN "clásicos" (definición + código + idea clave) usan el nuevo
  construir_slide_icn_flexible() que compone bloques dinámicamente en lugar
  de clonar un slide modelo rígido.
- Layouts especiales (anatomia, analogia, antes_despues, frase_clave) se
  mantienen sin cambios — ya funcionan bien.
- Modo --preview-icn para generar solo las slides de ICN (sin bienvenida,
  objetivo, haz ahora ni errores) para validación rápida de una sección.

Uso:
    python crear_ppt.py <ruta_spec.md> <ruta_salida.pptx>
    python crear_ppt.py <ruta_spec.md> <ruta_salida.pptx> --preview-icn

Requiere: pip install python-pptx
"""

import sys
import re
import copy
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Pt, Inches, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_ANCHOR
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("ERROR: python-pptx no está instalado. Instálalo con:")
    print("    pip install python-pptx")
    sys.exit(1)


CARPETA_SKILL = Path(__file__).resolve().parent
RUTA_PLANTILLA = CARPETA_SKILL / "plantilla_marca.pptx"

# — Importar helpers visuales de construir_plantilla —
sys.path.insert(0, str(CARPETA_SKILL))
try:
    from construir_plantilla import (
        aplicar_fondo_oscuro as _fondo_oscuro,
        agregar_barra_superior as _barra_superior,
        agregar_franja_lateral as _franja_lateral,
        agregar_rectangulo as _rect,
        agregar_texto as _texto,
        C as _C, T as _T, E as _E,
        ANCHO_SLIDE as _W, ALTO_SLIDE as _H, MARGEN_X as _MX,
    )
    _HELPERS_OK = True
except ImportError as _e:
    _HELPERS_OK = False
    print(f"[warn] No se pudo importar construir_plantilla: {_e}", file=sys.stderr)

# — Importar planificador —
try:
    from planificar_slides import planificar_icn, planificar_haz_ahora, SlidePlan, BloquePlan
    _PLANIFICADOR_OK = True
except ImportError as _e:
    _PLANIFICADOR_OK = False
    print(f"[warn] No se pudo importar planificar_slides: {_e}", file=sys.stderr)

M_BIENVENIDA = 0
M_OBJETIVO = 1
M_STOCK_TITULO = 2
M_CONCEPTO = 3
M_CODIGO_RESULTADO = 4
M_TABLA = 5
M_EJERCICIO = 6
M_TICKET = 7
M_CIERRE = 8
M_CODIGO_RESULTADO_MULTI = 9  # apilado (layout estrella)
M_ANATOMIA = 10                # ICN: anatomía de una expresión
M_ANALOGIA = 11                # ICN: analogía vida real ↔ Python
M_ANTES_DESPUES = 12           # ICN: comparación antes/después
M_FRASE_CLAVE = 13             # ICN: frase clave grande sola

TOTAL_SLIDES_MODELO = 14

# Ancho útil (in) de las cajas del slide apilado, para la regla mixta de fuente.
# Debe coincidir con COD_W / RES_W de construir_plantilla.py.
ANCHO_CAJA_CODIGO_MULTI = 6.85
ANCHO_CAJA_RESULTADO_MULTI = 13.333 - 0.6 - (0.6 + 6.85 + 0.30)

# Color verdoso para código inline (mismo que el de los bloques terminal)
COLOR_CODIGO_INLINE = RGBColor(0x4A, 0xDF, 0xCB)


# =====================================================================
# LOGGING DE DECISIONES DEL SELECTOR
# =====================================================================

import os
_DEBUG_PPT = os.environ.get("DEBUG_PPT") not in (None, "", "0", "false", "False")

def log_selector(mensaje: str):
    """Escribe una decisión del selector a stderr si DEBUG_PPT está activo.

    Activar con: DEBUG_PPT=1 python crear_ppt.py ...
    """
    if _DEBUG_PPT:
        print(f"[selector] {mensaje}", file=sys.stderr)

def log_warn(mensaje: str):
    """Warning siempre visible en stderr — riesgo de desborde u otros."""
    print(f"[warn] {mensaje}", file=sys.stderr)


# =====================================================================
# MEDICIÓN DE TEXTO (estimación de líneas y ajuste de fuente)
# =====================================================================

# Anchos de glifo promedio por punto, en unidades de "fracción del pt".
# Calibri es proporcional (más estrecho); Consolas es monoespaciada.
ANCHO_GLIFO_PROMEDIO = {
    "Calibri": 0.48,
    "Consolas": 0.60,
}


def estimar_lineas(texto: str, tamano_pt: int, ancho_in: float,
                   fuente: str = "Calibri") -> int:
    """Estima cuántas líneas ocupa `texto` en una caja con word_wrap.

    Heurística simple: cada glifo ocupa ANCHO_GLIFO_PROMEDIO[fuente] * tamano_pt.
    Cuenta los saltos de línea explícitos y suma el wrap de cada línea según
    cuántos caracteres caben por línea visual.
    """
    if not texto:
        return 0
    ratio = ANCHO_GLIFO_PROMEDIO.get(fuente, 0.50)
    ancho_util_pt = max((ancho_in - 0.30) * 72.0, 1.0)
    chars_por_linea = max(int(ancho_util_pt / (ratio * tamano_pt)), 10)
    total = 0
    for linea in texto.split("\n"):
        if not linea:
            total += 1
            continue
        total += max(1, -(-len(linea) // chars_por_linea))  # ceil div
    return total


# =====================================================================
# PARSER DEL SPEC
# =====================================================================

def parsear_spec(ruta_spec: Path) -> dict:
    contenido = ruta_spec.read_text(encoding="utf-8")
    spec = {
        "numero_clase": None, "tema": None, "contexto": {},
        "objetivo": None, "proposito": None,
        "haz_ahora": None,
        "icn": {"conceptos": [], "errores_tabla": [], "demos_apiladas": []},
        "guiada": {"situacion": None, "variables": None, "resultado": None},
        "independiente": [],
        "ticket": {"tarea": None, "entrega": None},
        "cierre": [],
    }
    m = re.search(r"^#\s+Clase\s+(\w+)\s*[—\-–]\s*(.+)$", contenido, re.MULTILINE)
    if m:
        raw = m.group(1)
        spec["numero_clase"] = int(raw) if raw.isdigit() else raw
        spec["tema"] = m.group(2).strip()
    seccion_ctx = extraer(contenido, "## Contexto", "## Objetivo")
    if seccion_ctx:
        for linea in seccion_ctx.split("\n"):
            mkv = re.match(r"^-\s+\*\*([^:]+):\*\*\s+(.+)$", linea.strip())
            if mkv:
                spec["contexto"][mkv.group(1).strip()] = mkv.group(2).strip()
    spec["objetivo"] = limpiar(extraer(contenido, "## Objetivo", "## Propósito"))
    prop = limpiar(extraer(contenido, "## Propósito", "## Estructura de la clase"))
    if prop:
        spec["proposito"] = "\n".join(
            l.lstrip("> ").rstrip() for l in prop.split("\n")
        ).strip()
    spec["haz_ahora"] = limpiar(extraer(contenido, "### 1. Haz Ahora", "### 2."))
    icn_raw = limpiar(extraer(contenido, "### 2. Introducción al Contenido Nuevo", "### 3."))
    if icn_raw:
        spec["icn"] = parsear_icn_enriquecido(icn_raw)
    g_raw = limpiar(extraer(contenido, "### 3. Práctica Guiada", "### 4."))
    if g_raw:
        spec["guiada"] = parsear_guiada(g_raw)
    i_raw = limpiar(extraer(contenido, "### 4. Práctica Independiente", "### 5."))
    if i_raw:
        spec["independiente"] = parsear_independiente(i_raw)
    t_raw = limpiar(extraer(contenido, "### 5. Ticket de Salida", "### Cierre"))
    if t_raw:
        m_tarea = re.search(r"\*\*Tarea:\*\*\s*(.+?)(?=\*\*Entrega|\Z)", t_raw, re.DOTALL)
        if m_tarea:
            spec["ticket"]["tarea"] = m_tarea.group(1).strip().strip("*").strip()
        m_ent = re.search(r"\*\*Entrega:\*\*\s*(.+)", t_raw)
        if m_ent:
            spec["ticket"]["entrega"] = m_ent.group(1).strip()
    c_raw = limpiar(extraer(contenido, "### Cierre y metacognición", "## Decisiones"))
    if c_raw:
        spec["cierre"] = [
            re.sub(r"^\d+\.\s*", "", l.strip())
            for l in c_raw.split("\n")
            if re.match(r"^\d+\.", l.strip())
        ]
    return spec


def extraer(texto: str, ini: str, fin):
    i = texto.find(ini)
    if i == -1:
        return None
    i += len(ini)
    if fin:
        f = texto.find(fin, i)
        return texto[i:] if f == -1 else texto[i:f]
    return texto[i:]


def limpiar(s):
    if not s:
        return None
    s = re.sub(r"\n\s*---\s*\n", "\n", s)
    return s.strip()


def parsear_icn_enriquecido(texto: str) -> dict:
    icn = {"conceptos": [], "errores_tabla": [], "demos_apiladas": []}
    bloques = re.split(r"\*\*Concepto\s+(\d+)\s*:\s*([^*]+)\*\*", texto)
    if len(bloques) > 1:
        for i in range(1, len(bloques), 3):
            if i + 2 >= len(bloques):
                break
            num = bloques[i].strip()
            nombre = bloques[i + 1].strip()
            cuerpo = bloques[i + 2]
            concepto = {
                "numero": num, "nombre": nombre,
                "definicion": "", "ejemplo": "", "idea_clave": "",
                # Campos del v5 (Etapa 1): override y subsecciones opcionales
                "tipo": None,            # "anatomia" | "analogia" | "antes_despues" | "frase_clave" | None
                "expresion": "",        # para anatomia: la expresión grande
                "partes": [],            # para anatomia: lista de {label, codigo, desc}
                "analogia_subtitulo": "",
                "analogia_filas": [],    # lista de {vida_real, python}
                "antes_label": "Antes",
                "antes_codigo": "",
                "despues_label": "Después",
                "despues_codigo": "",
                "takeaway": "",
                "frase_clave_aparte": False,
            }

            # Tipo: explícito (línea suelta o como bullet)
            m = re.search(r"(?:^|\n)\s*(?:[-•]\s*)?Tipo\s*:\s*([a-z_]+)",
                          cuerpo, re.IGNORECASE)
            if m:
                concepto["tipo"] = m.group(1).strip().lower()

            m = re.search(r"[-•]\s*\*?\*?Definici[óo]n\*?\*?\s*:\s*(.+?)(?=\n\s*[-•]|\n\s*\*\*|\Z)",
                          cuerpo, re.DOTALL)
            if m:
                concepto["definicion"] = m.group(1).strip()
            m = re.search(r"[-•]\s*\*?\*?Ejemplo\*?\*?\s*:?\s*\n?\s*```python\n(.*?)```",
                          cuerpo, re.DOTALL)
            if m:
                concepto["ejemplo"] = m.group(1).rstrip()
            m = re.search(r"[-•]\s*\*?\*?Idea clave\*?\*?\s*:\s*(.+?)(?=\n\s*[-•]|\n\s*\*\*|\Z)",
                          cuerpo, re.DOTALL)
            if m:
                concepto["idea_clave"] = m.group(1).strip()

            # Anatomía: **Expresión:** + **Partes:** (lista)
            # Soporta tanto **Expresión**: (colon fuera del bold) como **Expresión:** (colon dentro)
            m = re.search(r"\*\*Expresi[óo]n(?:\*\*)?\s*:(?:\*\*)?\s*`?([^`\n]+)`?",
                          cuerpo)
            if m:
                concepto["expresion"] = m.group(1).strip().strip("`").strip("*").strip()
            m_partes = re.search(
                r"\*\*Partes(?:\*\*)?\s*:(?:\*\*)?\s*\n(.*?)(?=\n\s*\*\*|\Z)",
                cuerpo, re.DOTALL)
            if m_partes:
                for linea in m_partes.group(1).split("\n"):
                    linea = linea.strip()
                    if not linea.startswith(("-", "•")):
                        continue
                    # Saltar líneas de "Idea clave:" que queden dentro del bloque
                    if re.match(r"[-•]\s*\*?\*?Idea\s+clave", linea, re.IGNORECASE):
                        continue
                    # Formato: - `código` — descripción [| label]
                    contenido = linea.lstrip("-•").strip()
                    label, codigo, desc = "", "", ""
                    m_cb = re.match(r"`([^`]+)`\s*[—\-:]\s*(.+)", contenido)
                    if m_cb:
                        codigo = m_cb.group(1).strip()
                        desc = m_cb.group(2).strip()
                    elif "|" in contenido:
                        partes = [p.strip() for p in contenido.split("|")]
                        if len(partes) == 3:
                            label, codigo, desc = partes
                        elif len(partes) == 2:
                            codigo, desc = partes
                    else:
                        desc = contenido
                    # Label puede venir como prefijo "**Nombre:** ..."
                    m_lbl = re.match(r"\*\*([^*]+)\*\*\s*:\s*(.+)", desc)
                    if m_lbl and not label:
                        label = m_lbl.group(1).strip()
                        desc = m_lbl.group(2).strip()
                    concepto["partes"].append({
                        "label": label, "codigo": codigo, "desc": desc,
                    })

            # Analogía: **Analogía:** (subtítulo opcional) + filas
            # Soporta **Analogía**: y **Analogía:** (colon dentro o fuera del bold)
            m_an = re.search(
                r"\*\*Analog[íi]a(?:\*\*)?\s*:(?:\*\*)?\s*([^\n]*)\n(.*?)(?=\n\s*\*\*|\Z)",
                cuerpo, re.DOTALL)
            if m_an:
                concepto["analogia_subtitulo"] = m_an.group(1).strip().strip("*").strip()
                for linea in m_an.group(2).split("\n"):
                    linea = linea.strip()
                    if not linea.startswith(("-", "•")):
                        continue
                    contenido = linea.lstrip("-•").strip()
                    # Formato: - vida real | python
                    if "|" in contenido:
                        vr, py = [p.strip() for p in contenido.split("|", 1)]
                        concepto["analogia_filas"].append({
                            "vida_real": vr, "python": py,
                        })
                    elif "↔" in contenido:
                        vr, py = [p.strip() for p in contenido.split("↔", 1)]
                        concepto["analogia_filas"].append({
                            "vida_real": vr, "python": py,
                        })

            # Antes/Después: **Antes:** ```python``` + **Después:** ```python```
            m_a = re.search(
                r"\*\*Antes\*?\*?(?:\s*\(([^)]+)\))?\s*:\s*\n```python\n(.*?)```",
                cuerpo, re.DOTALL)
            if m_a:
                if m_a.group(1):
                    concepto["antes_label"] = m_a.group(1).strip()
                concepto["antes_codigo"] = m_a.group(2).rstrip()
            m_d = re.search(
                r"\*\*Despu[ée]s\*?\*?(?:\s*\(([^)]+)\))?\s*:\s*\n```python\n(.*?)```",
                cuerpo, re.DOTALL)
            if m_d:
                if m_d.group(1):
                    concepto["despues_label"] = m_d.group(1).strip()
                concepto["despues_codigo"] = m_d.group(2).rstrip()
            m_tk = re.search(r"\*\*Takeaway\*?\*?\s*:\s*(.+?)(?=\n\s*\*\*|\Z)",
                             cuerpo, re.DOTALL)
            if m_tk:
                concepto["takeaway"] = m_tk.group(1).strip()

            # Flag: imprimir idea clave como slide frase_clave aparte
            m_fc = re.search(
                r"Frase clave aparte\s*:?\s*\*?\*?\s*(s[íi]|yes|true)",
                cuerpo, re.IGNORECASE)
            if m_fc:
                concepto["frase_clave_aparte"] = True

            # Demos que aparecen dentro del cuerpo de este concepto
            concepto["demos_inline"] = parsear_demos_apiladas(cuerpo)

            icn["conceptos"].append(concepto)

    if not icn["conceptos"]:
        m = re.search(r"\*\*Conceptos:\*\*\s*\n(.*?)(?=\*\*|\Z)", texto, re.DOTALL)
        if m:
            for j, l in enumerate(m.group(1).split("\n"), start=1):
                l = l.strip()
                if re.match(r"^\d+\.", l):
                    icn["conceptos"].append({
                        "numero": str(j), "nombre": f"Concepto {j}",
                        "definicion": re.sub(r"^\d+\.\s*", "", l),
                        "ejemplo": "", "idea_clave": "",
                    })
        m = re.search(r"\*\*Ejemplos a usar:\*\*\s*\n```python\n(.*?)```",
                      texto, re.DOTALL)
        if m and icn["conceptos"]:
            lineas = [l for l in m.group(1).split("\n") if l.strip()]
            for c, lin in zip(icn["conceptos"], lineas):
                c["ejemplo"] = lin

    m = re.search(r"\*\*Errores típicos:\*\*\s*\n(\|.*?)(?=\n\n|\Z)", texto, re.DOTALL)
    if m:
        filas = []
        for linea in m.group(1).strip().split("\n"):
            linea = linea.strip()
            if not linea.startswith("|"):
                continue
            celdas = [c.strip() for c in linea.split("|")[1:-1]]
            if all(re.match(r"^[-:]+$", c) for c in celdas):
                continue
            if celdas:
                filas.append(celdas)
        if filas:
            filas = filas[1:]
        icn["errores_tabla"] = filas

    # Demostraciones apiladas (slide tipo 9). Sintaxis opcional en el spec:
    #
    #   **Demostración: [título corto del slide]**
    #   - Fila: [etiqueta] | `código en una o varias líneas con \n` | resultado
    #   ... (hasta 3 filas)
    #
    # También admite código multilínea con bloque ```python```:
    #   - Fila: [etiqueta]
    #     ```python
    #     línea1
    #     línea2
    #     ```
    #     Resultado: salida esperada
    icn["demos_apiladas"] = parsear_demos_apiladas(texto)
    return icn


def parsear_demos_apiladas(texto: str) -> list:
    """Parsea bloques '**Demostración: ...**' en una lista de demos.

    Cada demo = {"titulo": str, "subtitulo": str, "filas": [ {label, codigo, resultado} ]}
    Soporta dos formatos de fila:
      A) en línea:  - Fila: etiqueta | codigo | resultado    (codigo admite \\n literal)
      B) con bloque de código:
         - Fila: etiqueta
           ```python
           ...
           ```
           Resultado: ...
    """
    demos = []
    bloques = re.split(r"\*\*Demostración\s*:\s*([^*]+)\*\*", texto)
    if len(bloques) < 2:
        return demos
    for i in range(1, len(bloques), 2):
        titulo = bloques[i].strip()
        cuerpo = bloques[i + 1] if i + 1 < len(bloques) else ""
        demo = {"titulo": titulo, "subtitulo": "", "filas": []}

        # Subtítulo opcional: primera línea que empieza con "Subtítulo:"
        m_sub = re.search(r"Subt[íi]tulo\s*:\s*(.+)", cuerpo)
        if m_sub:
            demo["subtitulo"] = m_sub.group(1).strip()

        # Formato B: filas con bloque ```python```
        patron_b = re.compile(
            r"[-•]\s*Fila\s*:\s*([^\n]*?)\s*\n\s*```python\s*\n(.*?)```"
            r"(?:\s*Resultado\s*:\s*([^\n]+))?",
            re.DOTALL)
        encontrados_b = list(patron_b.finditer(cuerpo))
        if encontrados_b:
            for mb in encontrados_b:
                label = mb.group(1).strip()
                codigo = mb.group(2).rstrip()
                resultado = (mb.group(3) or "").strip()
                demo["filas"].append({"label": label, "codigo": codigo,
                                      "resultado": resultado})
        else:
            # Formato A: una línea por fila con separador |
            for linea in cuerpo.split("\n"):
                linea = linea.strip()
                m_fila = re.match(r"[-•]\s*Fila\s*:\s*(.+)", linea)
                if not m_fila:
                    continue
                partes = [p.strip() for p in m_fila.group(1).split("|")]
                label = partes[0] if len(partes) > 0 else ""
                codigo = partes[1].replace("\\n", "\n") if len(partes) > 1 else ""
                resultado = partes[2].replace("\\n", "\n") if len(partes) > 2 else ""
                demo["filas"].append({"label": label, "codigo": codigo,
                                      "resultado": resultado})

        if demo["filas"]:
            demos.append(demo)
    return demos


def parsear_guiada(texto: str) -> dict:
    g = {"situacion": None, "variables": None, "resultado": None}
    m = re.search(r"\*\*Situación:\*\*\s*(.+?)(?=\*\*|\Z)", texto, re.DOTALL)
    if m:
        g["situacion"] = m.group(1).strip()
    m = re.search(r"\*\*Variables:\*\*\s*\n```python\n(.*?)```", texto, re.DOTALL)
    if m:
        g["variables"] = m.group(1).rstrip()
    m = re.search(r"\*\*Resultado esperado:\*\*\s*\n```\n?(.*?)```", texto, re.DOTALL)
    if m:
        g["resultado"] = m.group(1).rstrip()
    return g


def parsear_independiente(texto: str) -> list:
    ejercicios = []
    bloques = re.split(r"\*\*Ejercicio\s+(\d+)\s*[—\-–]\s*([^*]+)\*\*", texto)
    for i in range(1, len(bloques), 3):
        if i + 2 >= len(bloques):
            break
        num = bloques[i].strip()
        titulo = bloques[i + 1].strip()
        cuerpo = bloques[i + 2].strip()
        ej = {"numero": num, "titulo": titulo, "enunciado": None, "resultado": None}
        m = re.search(r"Resultado esperado:\s*\n```\n?(.*?)```", cuerpo, re.DOTALL)
        if m:
            ej["resultado"] = m.group(1).rstrip()
            ej["enunciado"] = cuerpo[:m.start()].strip()
        else:
            ej["enunciado"] = cuerpo
        ejercicios.append(ej)
    return ejercicios


def separar_haz_ahora(texto: str) -> dict:
    """Separa el Haz Ahora en texto narrativo + bloque de código (si existe)."""
    partes = {"intro": "", "texto": "", "codigo": ""}
    if not texto:
        return partes

    # Quitar "**Propósito:**" si está
    texto = re.sub(r"\*\*Propósito:\*\*[^\n]*\n?", "", texto)
    # Quitar línea sola con "(N min)" que es metadata
    texto = re.sub(r"^\s*\(\s*\d+\s*min\s*\)\s*$", "", texto, flags=re.MULTILINE)

    # Capturar intro desde "**Actividad:**" hasta el final del párrafo
    m_act = re.search(r"\*\*Actividad:\*\*\s*(.+?)(?=\n|\Z)", texto)
    if m_act:
        partes["intro"] = m_act.group(1).strip()

    # Capturar el bloque de código markdown
    m_code = re.search(r"```python\s*\n(.*?)```", texto, re.DOTALL)
    if m_code:
        partes["codigo"] = m_code.group(1).rstrip()
        sin_codigo = re.sub(r"```python\s*\n.*?```", "", texto, flags=re.DOTALL)
        sin_codigo = re.sub(r"\*\*Actividad:\*\*[^\n]*\n?", "", sin_codigo)
        # Compactar líneas en blanco múltiples
        sin_codigo = re.sub(r"\n\s*\n", "\n\n", sin_codigo)
        partes["texto"] = sin_codigo.strip()
    else:
        sin_actividad = re.sub(r"\*\*Actividad:\*\*[^\n]*\n?", "", texto)
        sin_actividad = re.sub(r"\n\s*\n", "\n\n", sin_actividad)
        partes["texto"] = sin_actividad.strip()

    if not partes["intro"]:
        partes["intro"] = "Activa lo que ya sabes antes de empezar."
    if not partes["texto"]:
        partes["texto"] = "Observa con atención y prepárate para discutir."
    return partes


# =====================================================================
# DUPLICACIÓN Y REEMPLAZO
# =====================================================================

def duplicar_slide(prs, slide_modelo):
    blank_layout = slide_modelo.slide_layout
    nueva = prs.slides.add_slide(blank_layout)
    for shp in list(nueva.shapes):
        sp = shp._element
        sp.getparent().remove(sp)
    for shp in slide_modelo.shapes:
        nuevo_el = copy.deepcopy(shp._element)
        nueva.shapes._spTree.insert_element_before(nuevo_el, "p:extLst")
    return nueva


def eliminar_slides(prs, indices):
    xml_slides = prs.slides._sldIdLst
    slides_xml = list(xml_slides)
    for idx in sorted(indices, reverse=True):
        rId = slides_xml[idx].rId
        prs.part.drop_rel(rId)
        xml_slides.remove(slides_xml[idx])


def reemplazar_placeholders(slide, mapping: dict, mapping_inline_code=None):
    """Reemplaza placeholders {{KEY}} en todos los textos y tablas.

    Si el valor del placeholder contiene backticks `código`, automáticamente
    se aplica formato Consolas verdoso a esa parte. Esto aplica a TODOS los
    placeholders por defecto. El parámetro mapping_inline_code se mantiene
    por compatibilidad pero ya no es necesario.
    """
    for shape in slide.shapes:
        if shape.has_text_frame:
            reemplazar_en_text_frame(shape.text_frame, mapping)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    reemplazar_en_text_frame(cell.text_frame, mapping)


def reemplazar_en_text_frame(tf, mapping: dict):
    """Para cada run del TextFrame, busca placeholders y los reemplaza.

    Si el valor de un placeholder contiene backticks, el texto se reconstruye
    en múltiples runs para aplicar formato Consolas verdoso al código inline.
    Si el valor contiene saltos de línea \\n, se crean párrafos adicionales.
    """
    for para in tf.paragraphs:
        for run in list(para.runs):
            texto_original = run.text
            placeholders_encontrados = re.findall(r"\{\{([A-Z_0-9]+)\}\}", texto_original)
            if not placeholders_encontrados:
                continue

            # ¿Algún placeholder tiene backticks en su valor?
            necesita_split_backticks = any(
                "`" in str(mapping.get(k, "")) for k in placeholders_encontrados
            )
            # ¿Algún placeholder tiene saltos de línea?
            necesita_split_lineas = any(
                "\n" in str(mapping.get(k, "")) for k in placeholders_encontrados
            )

            if not necesita_split_backticks and not necesita_split_lineas:
                nuevo = texto_original
                for key, valor in mapping.items():
                    nuevo = nuevo.replace("{{" + key + "}}",
                                          str(valor) if valor is not None else "")
                run.text = nuevo
            elif necesita_split_lineas and not necesita_split_backticks:
                aplicar_reemplazo_con_saltos_linea(tf, para, run, mapping)
            else:
                # Tiene backticks (con o sin saltos de línea)
                aplicar_reemplazo_con_inline_code(para, run, mapping)


def aplicar_reemplazo_con_saltos_linea(tf, para, run, mapping):
    """Maneja placeholders cuyo valor contiene saltos de línea.

    En XML de PowerPoint, los párrafos son elementos <a:p>. Un \\n debe
    convertirse en un nuevo elemento <a:p> hermano. Esta función reescribe
    el run actual para que su texto sea solo la primera línea, y agrega
    párrafos adicionales después con el resto del contenido.
    """
    texto_original = run.text

    # Conservar formato base del run original
    base_font_name = run.font.name
    base_font_size = run.font.size
    base_font_bold = run.font.bold
    base_color = None
    try:
        if run.font.color and run.font.color.type:
            base_color = run.font.color.rgb
    except Exception:
        pass

    # Expandir placeholders
    texto_expandido = texto_original
    for key, valor in mapping.items():
        texto_expandido = texto_expandido.replace(
            "{{" + key + "}}", str(valor) if valor is not None else ""
        )

    lineas = texto_expandido.split("\n")
    if not lineas:
        return

    # Primera línea va en el run actual
    run.text = lineas[0]

    # Líneas adicionales como párrafos nuevos después del párrafo actual
    if len(lineas) > 1:
        from copy import deepcopy
        p_xml = para._p
        parent = p_xml.getparent()
        idx_actual = list(parent).index(p_xml)

        for offset, linea in enumerate(lineas[1:], start=1):
            # Clonar el párrafo actual como base (mantiene formato)
            nuevo_p = deepcopy(p_xml)
            # Limpiar runs dentro del nuevo párrafo
            ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
            for r_xml in list(nuevo_p.findall(f"{ns}r")):
                nuevo_p.remove(r_xml)
            # Insertar el nuevo párrafo después del actual
            parent.insert(idx_actual + offset, nuevo_p)

        # Ahora poblar las líneas en los nuevos párrafos
        # Re-leer paragraphs después de la mutación XML
        nuevos_paras = list(tf.paragraphs)
        idx_para_actual = None
        for i, p in enumerate(nuevos_paras):
            if p._p is p_xml:
                idx_para_actual = i
                break
        if idx_para_actual is not None:
            for offset, linea in enumerate(lineas[1:], start=1):
                p_nuevo = nuevos_paras[idx_para_actual + offset]
                r_nuevo = p_nuevo.add_run()
                r_nuevo.text = linea
                r_nuevo.font.name = base_font_name or "Calibri"
                if base_font_size:
                    r_nuevo.font.size = base_font_size
                r_nuevo.font.bold = base_font_bold
                if base_color:
                    r_nuevo.font.color.rgb = base_color


def aplicar_reemplazo_con_inline_code(para, run, mapping):
    """Toma un run que contiene placeholders con backticks y lo reemplaza por
    múltiples runs: texto normal y código (Consolas verdoso) alternados."""
    texto_original = run.text

    # Conservar el formato del run original como base
    base_font_name = run.font.name
    base_font_size = run.font.size
    base_font_bold = run.font.bold
    base_font_color_rgb = None
    try:
        if run.font.color and run.font.color.type:
            base_font_color_rgb = run.font.color.rgb
    except Exception:
        pass

    # Primero, sustituir placeholders por su valor literal
    texto_expandido = texto_original
    for key, valor in mapping.items():
        texto_expandido = texto_expandido.replace(
            "{{" + key + "}}", str(valor) if valor is not None else ""
        )

    # Limpiar el run original (lo vamos a reescribir como secuencia de runs)
    run.text = ""

    # Dividir por backticks: alternar texto plano y código
    partes = re.split(r"(`[^`]+`)", texto_expandido)

    # El primer run ya existe; rellenamos su texto con la primera parte
    primera = True
    para_xml = para._p
    # Limpiar a partir del run original todos los runs siguientes que estaban en este párrafo
    # (estamos dentro del bucle paraventana sobre runs, pero como ya consumimos los
    # placeholders del original solo, esto está OK)

    for parte in partes:
        if not parte:
            continue
        if parte.startswith("`") and parte.endswith("`"):
            # Código inline
            contenido = parte[1:-1]
            if primera:
                r = run
                primera = False
            else:
                r = para.add_run()
            r.text = contenido
            r.font.name = "Consolas"
            if base_font_size:
                r.font.size = base_font_size
            r.font.bold = base_font_bold
            r.font.color.rgb = COLOR_CODIGO_INLINE
        else:
            # Texto normal
            if primera:
                r = run
                primera = False
            else:
                r = para.add_run()
            r.text = parte
            r.font.name = base_font_name or "Calibri"
            if base_font_size:
                r.font.size = base_font_size
            r.font.bold = base_font_bold
            if base_font_color_rgb:
                r.font.color.rgb = base_font_color_rgb


# =====================================================================
# REGLA MIXTA DE FUENTE PARA CÓDIGO (ajuste sobre contenido real)
# =====================================================================

def calcular_tamano_codigo(texto, ancho_caja_in, tam_max=20, tam_min=16):
    """Elige el mayor tamaño (20 -> 18 -> 16) con que la línea más larga cabe.

    Consolas es monoespaciada (~0.60 * pt de ancho por glifo). ancho_caja_in
    es el ancho útil de la caja en pulgadas. Devuelve un entero de tamaño en pt.
    """
    lineas = (texto or "").split("\n")
    max_chars = max((len(l) for l in lineas), default=0)
    if max_chars == 0:
        return tam_max
    ancho_util_pt = (ancho_caja_in - 0.30) * 72.0  # ~0.15" de padding por lado
    for tam in sorted({tam_max, 18, tam_min}, reverse=True):
        if max_chars * (0.60 * tam) <= ancho_util_pt:
            return tam
    return tam_min


def ajustar_fuente_codigo_en_shape(slide, placeholder_key, ancho_caja_in,
                                   texto_real, tam_max=20, tam_min=16):
    """Tras reemplazar un placeholder de código, reajusta el tamaño de fuente
    de los runs Consolas del shape que contenía ese placeholder.

    Se llama ANTES del reemplazo: localiza el shape cuyo text_frame todavía
    contiene "{{KEY}}", calcula el tamaño según texto_real y lo aplica a todos
    sus runs. Devuelve el tamaño aplicado (o None si no encontró el shape).
    """
    tam = calcular_tamano_codigo(texto_real, ancho_caja_in, tam_max, tam_min)
    objetivo = "{{" + placeholder_key + "}}"
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if objetivo in shape.text_frame.text:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(tam)
            return tam
    return None


def buscar_shape_por_key(slide, key):
    """Encuentra la shape que contiene {{KEY}} en su texto (antes de reemplazar)."""
    token = "{{" + key + "}}"
    for shape in slide.shapes:
        if shape.has_text_frame:
            texto = " ".join(r.text for p in shape.text_frame.paragraphs for r in p.runs)
            if token in texto:
                return shape
    return None


def eliminar_shape_de_slide(slide, shape):
    """Elimina una shape del XML del slide."""
    sp = shape._element
    sp.getparent().remove(sp)


def extraer_tabla_md(texto):
    """Extrae una tabla markdown del texto.
    Devuelve (texto_sin_tabla, lista_de_filas) o (texto_original, None).
    Cada fila es una lista de strings (celdas). La fila separadora (---) se omite.
    """
    if not texto or "|" not in texto:
        return texto, None
    patron = r"((?:\|.+\|\n?)+)"
    m = re.search(patron, texto)
    if not m:
        return texto, None
    tabla_texto = m.group(1)
    texto_limpio = (texto[:m.start()] + texto[m.end():]).strip()
    filas = []
    for linea in tabla_texto.strip().split("\n"):
        if re.match(r"^\|[\s\-|]+\|$", linea.strip()):
            continue  # fila separadora
        celdas = [c.strip() for c in linea.strip().strip("|").split("|")]
        if any(c for c in celdas):
            filas.append(celdas)
    return texto_limpio, filas if filas else None


def insertar_tabla_pptx(slide, filas, left, top, width, height):
    """Agrega una tabla estilizada (marca oscura) al slide en la posición dada."""
    from pptx.util import Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    if not filas:
        return
    n_rows = len(filas)
    n_cols = max(len(f) for f in filas)

    tbl = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table

    COLOR_HEADER_BG = RGBColor(0x0D, 0x2B, 0x3E)
    COLOR_ROW_BG    = RGBColor(0x12, 0x22, 0x32)
    COLOR_ROW_ALT   = RGBColor(0x0A, 0x1A, 0x28)
    COLOR_TEXTO     = RGBColor(0xFF, 0xFF, 0xFF)
    COLOR_CODIGO    = RGBColor(0x4A, 0xDF, 0xCB)

    for r_idx, fila in enumerate(filas):
        for c_idx in range(n_cols):
            cell = tbl.cell(r_idx, c_idx)
            cell.margin_left  = Emu(91440)
            cell.margin_right = Emu(91440)
            cell.margin_top   = Emu(45720)
            cell.margin_bottom= Emu(45720)
            cell.fill.solid()
            if r_idx == 0:
                cell.fill.fore_color.rgb = COLOR_HEADER_BG
            elif r_idx % 2 == 1:
                cell.fill.fore_color.rgb = COLOR_ROW_BG
            else:
                cell.fill.fore_color.rgb = COLOR_ROW_ALT

            texto_celda = fila[c_idx] if c_idx < len(fila) else ""
            tf = cell.text_frame
            tf.word_wrap = True
            para = tf.paragraphs[0]
            # Detectar si hay backticks (código inline)
            partes = re.split(r"`([^`]+)`", texto_celda)
            for k, parte in enumerate(partes):
                if not parte:
                    continue
                run = para.add_run()
                run.text = parte
                run.font.size = Pt(14) if r_idx == 0 else Pt(13)
                run.font.bold = r_idx == 0
                if k % 2 == 1:  # estaba entre backticks
                    run.font.name = "Consolas"
                    run.font.color.rgb = COLOR_CODIGO
                else:
                    run.font.color.rgb = COLOR_TEXTO


def adaptar_layout_concepto(slide, concepto):
    """Post-procesa un slide de concepto para adaptar el layout al contenido real.
    Debe llamarse DESPUÉS de reemplazar_placeholders.
    Recibe las shape references guardadas antes del reemplazo.
    """
    pass  # lógica delegada a las referencias guardadas en generar_ppt


def ocultar_filas_vacias_multi(slide, filas_vacias):
    """Elimina las cajas, etiquetas y textos de las filas no usadas del slide
    apilado, identificándolas por su rango vertical EXACTO (no por proximidad).

    filas_vacias: lista de números de fila (1..3) que deben ocultarse.

    La geometría debe coincidir con construir_slide_codigo_resultado_multi:
      FILA_Y0 = 2.25, FILA_PASO = 1.62, LABEL_H = 0.34, FILA_H = 1.42
    Una fila i ocupa verticalmente desde y_label hasta el fondo de la caja.
    """
    FILA_Y0 = 2.25
    FILA_PASO = 1.62
    LABEL_H = 0.34
    FILA_H = 1.42
    MARGEN = 0.05  # tolerancia

    rangos = []
    for i in filas_vacias:
        y_top = FILA_Y0 + (i - 1) * FILA_PASO - MARGEN
        y_bottom = y_top + LABEL_H + FILA_H + 2 * MARGEN
        rangos.append((Inches(y_top), Inches(y_bottom)))

    for shape in list(slide.shapes):
        try:
            top = shape.top
            bottom = shape.top + shape.height
        except Exception:
            continue
        for (y0, y1) in rangos:
            # El shape pertenece a la fila si su centro cae dentro del rango
            centro = (top + bottom) // 2
            if y0 <= centro <= y1:
                shape._element.getparent().remove(shape._element)
                break


# =====================================================================
# SELECTOR DE LAYOUT — reglas §7 del plan
# =====================================================================

# Umbrales para detección automática de tipo ICN
UMBRAL_DEFINICION_CHARS = 220  # > esto sugiere usar anatomía/split
UMBRAL_IDEA_CLAVE_CHARS = 120  # > esto considera frase aparte


def seleccionar_layout_concepto(concepto: dict) -> str:
    """Decide qué layout usar para un bloque de concepto ICN.

    Devuelve uno de: "concepto" | "anatomia" | "analogia" | "antes_despues"
                    | "frase_clave"

    Regla de prioridad:
      1. Tipo: explícito en el spec.
      2. Si el bloque tiene Partes parseadas → anatomia.
      3. Si el bloque tiene filas de analogía → analogia.
      4. Si el bloque tiene antes_codigo y despues_codigo → antes_despues.
      5. Default → concepto (layout actual con def + ejemplo + idea).
    """
    tipo_explicito = (concepto.get("tipo") or "").lower()
    if tipo_explicito in {"anatomia", "anatomía", "anatomy"}:
        return "anatomia"
    if tipo_explicito in {"analogia", "analogía", "analogy"}:
        return "analogia"
    if tipo_explicito in {"antes_despues", "antes-despues", "comparison",
                          "comparacion", "comparación"}:
        return "antes_despues"
    if tipo_explicito in {"frase_clave", "frase-clave", "pull_quote",
                          "pull-quote", "quote"}:
        return "frase_clave"
    if tipo_explicito in {"concepto", "definition", "default"}:
        return "concepto"

    # Inferencia automática
    if concepto.get("partes"):
        return "anatomia"
    if concepto.get("analogia_filas"):
        return "analogia"
    if concepto.get("antes_codigo") and concepto.get("despues_codigo"):
        return "antes_despues"

    # Aviso si la definición es muy larga y no hay tipo explícito
    if len(concepto.get("definicion", "")) > UMBRAL_DEFINICION_CHARS:
        log_warn(
            f"concepto {concepto.get('numero')} definición con "
            f"{len(concepto['definicion'])} chars > {UMBRAL_DEFINICION_CHARS}"
            f" — considera Tipo: anatomia o split"
        )
    return "concepto"


# =====================================================================
# CONSTRUCTORES DE SLIDES ICN
# =====================================================================

def construir_concepto_clasico(prs, num: str, concepto: dict):
    """Layout original: Definición + Ejemplo + Idea clave (con adaptaciones)."""
    sl = duplicar_slide(prs, prs.slides[M_CONCEPTO])

    tiene_codigo = bool(concepto["ejemplo"])
    texto_def = concepto["definicion"] or ""
    texto_sin_tabla, filas_tabla = extraer_tabla_md(texto_def)
    tiene_tabla = filas_tabla is not None

    shape_def  = buscar_shape_por_key(sl, "DEFINICION")
    shape_cod  = buscar_shape_por_key(sl, "EJEMPLO_CODIGO")
    shape_idea = buscar_shape_por_key(sl, "IDEA_CLAVE")

    reemplazar_placeholders(sl, {
        "NUMERO_CLASE": num,
        "SECCION": "Contenido nuevo",
        "TITULO_CONCEPTO": f"📘 {concepto['numero']}. {concepto['nombre']}",
        "DEFINICION": texto_sin_tabla if tiene_tabla else texto_def,
        "EJEMPLO_CODIGO": concepto["ejemplo"] or "",
        "IDEA_CLAVE": concepto["idea_clave"] or "",
    })

    if tiene_tabla and shape_cod:
        left, top = shape_cod.left, shape_cod.top
        width, height = shape_cod.width, shape_cod.height
        eliminar_shape_de_slide(sl, shape_cod)
        insertar_tabla_pptx(sl, filas_tabla, left, top, width, height)
    elif not tiene_codigo and shape_cod and shape_def:
        cod_height = shape_cod.height
        eliminar_shape_de_slide(sl, shape_cod)
        if shape_idea:
            shape_def.height  += cod_height * 2 // 3
            shape_idea.top    -= cod_height // 3
            shape_idea.height += cod_height // 3
        else:
            shape_def.height += cod_height

    # Sin idea clave: eliminar esa caja y dar el espacio al código
    if not concepto.get("idea_clave") and shape_idea and tiene_codigo and not tiene_tabla:
        idea_height = shape_idea.height
        eliminar_shape_de_slide(sl, shape_idea)
        if shape_cod:
            shape_cod.height += idea_height


def construir_concepto_anatomia(prs, num: str, concepto: dict):
    """Layout anatomía: expresión grande + hasta 4 partes en grilla 2x2."""
    sl = duplicar_slide(prs, prs.slides[M_ANATOMIA])

    partes = (concepto.get("partes") or [])[:4]
    expresion = concepto.get("expresion") or concepto.get("ejemplo") or ""
    # Si no hay expresión explícita pero hay código de ejemplo, usar primera línea
    if not expresion and concepto.get("ejemplo"):
        expresion = concepto["ejemplo"].split("\n", 1)[0]

    mapping = {
        "NUMERO_CLASE": num,
        "SECCION": "Contenido nuevo",
        "TITULO_SLIDE": f"📘 {concepto['numero']}. {concepto['nombre']}",
        "EXPRESION": expresion,
        "IDEA_CLAVE": concepto.get("idea_clave", ""),
    }
    for i in range(1, 5):
        if i <= len(partes):
            p = partes[i - 1]
            mapping[f"PARTE_{i}_LABEL"] = p.get("label", "") or f"PARTE {i}"
            mapping[f"PARTE_{i}_CODIGO"] = p.get("codigo", "")
            mapping[f"PARTE_{i}_DESC"] = p.get("desc", "")
        else:
            mapping[f"PARTE_{i}_LABEL"] = ""
            mapping[f"PARTE_{i}_CODIGO"] = ""
            mapping[f"PARTE_{i}_DESC"] = ""

    reemplazar_placeholders(sl, mapping)

    # Ocultar tarjetas no usadas (por geometría: las posiciones 3 y 4
    # ocupan la segunda fila de la grilla)
    if len(partes) < 4:
        _ocultar_tarjetas_anatomia(sl, partes_usadas=len(partes))


def _ocultar_tarjetas_anatomia(slide, partes_usadas: int):
    """Elimina las shapes de las tarjetas no usadas (3 y 4 viven en y>=4.6)."""
    # Geometría de construir_slide_anatomia
    EXP_Y = 1.70
    EXP_H = 0.95
    GRID_Y0 = EXP_Y + EXP_H + 0.25  # = 2.90
    CARD_H = 1.40
    CARD_GAP_Y = 0.20
    Y_FILA2 = GRID_Y0 + CARD_H + CARD_GAP_Y  # 4.50
    Y_FILA1 = GRID_Y0

    # Anchos de cada columna
    CARD_W = (13.333 - 2 * 0.6 - 0.30) / 2
    X_COL2 = 0.6 + CARD_W + 0.30

    # Rangos a borrar
    rangos = []  # cada uno: (x_min, x_max, y_min, y_max)
    if partes_usadas < 4:
        rangos.append((X_COL2 - 0.05, X_COL2 + CARD_W + 0.05,
                       Y_FILA2 - 0.05, Y_FILA2 + CARD_H + 0.05))
    if partes_usadas < 3:
        rangos.append((0.6 - 0.05, 0.6 + CARD_W + 0.05,
                       Y_FILA2 - 0.05, Y_FILA2 + CARD_H + 0.05))
    if partes_usadas < 2:
        rangos.append((X_COL2 - 0.05, X_COL2 + CARD_W + 0.05,
                       Y_FILA1 - 0.05, Y_FILA1 + CARD_H + 0.05))

    if not rangos:
        return
    for shape in list(slide.shapes):
        try:
            left_in = shape.left / 914400
            top_in = shape.top / 914400
            width_in = shape.width / 914400
            height_in = shape.height / 914400
        except Exception:
            continue
        cx = left_in + width_in / 2
        cy = top_in + height_in / 2
        for (x0, x1, y0, y1) in rangos:
            if x0 <= cx <= x1 and y0 <= cy <= y1:
                shape._element.getparent().remove(shape._element)
                break


def construir_concepto_analogia(prs, num: str, concepto: dict):
    """Layout analogía: hasta 4 filas vida real ↔ Python."""
    sl = duplicar_slide(prs, prs.slides[M_ANALOGIA])
    filas = (concepto.get("analogia_filas") or [])[:4]

    mapping = {
        "NUMERO_CLASE": num,
        "SECCION": "Contenido nuevo",
        "TITULO_SLIDE": f"📘 {concepto['numero']}. {concepto['nombre']}",
        "SUBTITULO_SLIDE": concepto.get("analogia_subtitulo")
                           or concepto.get("definicion", ""),
    }
    for i in range(1, 5):
        if i <= len(filas):
            f = filas[i - 1]
            mapping[f"FILA_{i}_VIDA_REAL"] = f.get("vida_real", "")
            mapping[f"FILA_{i}_PYTHON"] = f.get("python", "")
        else:
            mapping[f"FILA_{i}_VIDA_REAL"] = ""
            mapping[f"FILA_{i}_PYTHON"] = ""

    reemplazar_placeholders(sl, mapping)

    # Ocultar filas no usadas (por geometría: filas i ocupan y = FILA_Y0 + (i-1)*paso)
    if len(filas) < 4:
        _ocultar_filas_analogia(sl, filas_usadas=len(filas))


def _ocultar_filas_analogia(slide, filas_usadas: int):
    FILA_Y0 = 2.65
    FILA_H = 0.85
    FILA_GAP = 0.12
    rangos = []
    for i in range(filas_usadas + 1, 5):
        y = FILA_Y0 + (i - 1) * (FILA_H + FILA_GAP)
        rangos.append((y - 0.05, y + FILA_H + 0.05))
    if not rangos:
        return
    for shape in list(slide.shapes):
        try:
            top_in = shape.top / 914400
            height_in = shape.height / 914400
        except Exception:
            continue
        cy = top_in + height_in / 2
        for (y0, y1) in rangos:
            if y0 <= cy <= y1:
                shape._element.getparent().remove(shape._element)
                break


def construir_concepto_antes_despues(prs, num: str, concepto: dict):
    """Layout comparación antes/después con dos snippets paralelos."""
    sl = duplicar_slide(prs, prs.slides[M_ANTES_DESPUES])

    takeaway = (concepto.get("takeaway")
                or concepto.get("idea_clave")
                or "")
    subtitulo = concepto.get("definicion") or ""

    mapping = {
        "NUMERO_CLASE": num,
        "SECCION": "Contenido nuevo",
        "TITULO_SLIDE": f"📘 {concepto['numero']}. {concepto['nombre']}",
        "SUBTITULO_SLIDE": subtitulo,
        "ANTES_LABEL": (concepto.get("antes_label") or "Antes").upper(),
        "ANTES_CODIGO": concepto.get("antes_codigo", ""),
        "DESPUES_LABEL": (concepto.get("despues_label") or "Después").upper(),
        "DESPUES_CODIGO": concepto.get("despues_codigo", ""),
        "TAKEAWAY": takeaway,
    }
    reemplazar_placeholders(sl, mapping)


def construir_concepto_frase_clave(prs, num: str, concepto: dict,
                                   solo_frase: bool = False):
    """Layout frase clave grande sola. Útil como pausa visual.

    Si solo_frase=True, se usa la `idea_clave` como FRASE y `nombre` como
    ATRIBUCION (modo "pausa después del concepto"). Si solo_frase=False,
    se asume que el spec tiene Tipo: frase_clave directo: en ese caso la
    frase es la propia definición o la idea clave.
    """
    sl = duplicar_slide(prs, prs.slides[M_FRASE_CLAVE])

    if solo_frase:
        frase = concepto.get("idea_clave", "")
        atribucion = f"Concepto {concepto.get('numero','?')} · {concepto.get('nombre','')}"
    else:
        frase = (concepto.get("definicion")
                 or concepto.get("idea_clave")
                 or "")
        atribucion = concepto.get("nombre", "")

    reemplazar_placeholders(sl, {
        "NUMERO_CLASE": num,
        "SECCION": "Contenido nuevo",
        "FRASE": frase,
        "ATRIBUCION": atribucion,
    })


# =====================================================================
# SLIDE ICN FLEXIBLE — composición dinámica de bloques
# =====================================================================

# Geometría del área de contenido (debajo del título, sobre el margen inferior)
_Y_ICN_INICIO = 1.65   # pulgadas desde el tope del slide
_Y_ICN_FIN    = 7.20   # margen inferior
_GAP_BLOQUES  = 0.18   # espacio entre bloques

# Anchos útiles para medición de fuente de código
_ANCHO_COD_FLEX = _W - 2 * _MX - 0.3 if _HELPERS_OK else 12.0


def _renderizar_bloque(slide, bloque: "BloquePlan", y: float, h: float,
                       key: str, mapping: dict) -> None:
    """Dibuja un bloque en la posición y con altura h.

    Para el texto con backticks, usa el sistema de placeholders de
    reemplazar_placeholders (que formatea código inline automáticamente).
    """
    if not _HELPERS_OK:
        return

    tipo = bloque.tipo

    if tipo == "definicion":
        LABEL_H = 0.38
        _rect(slide, _MX, y, _W - 2 * _MX, h, _C["fondo_bloque"])
        _rect(slide, _MX, y, 0.08, h, _C["turquesa"])
        label_txt = f"📘  {bloque.label or 'Definición'}"
        _texto(slide, _MX + 0.25, y + 0.10, _W - 2 * _MX - 0.4, LABEL_H,
               label_txt, tamano=_T["titulo_bloque"], color="turquesa", negrita=True)
        # Tamaño de fuente adaptado a la longitud
        tam_def = 20 if len(str(bloque.contenido)) > 150 else 22
        _texto(slide, _MX + 0.25, y + LABEL_H + 0.05,
               _W - 2 * _MX - 0.4, h - LABEL_H - 0.15,
               "{{" + key + "_DEF}}", tamano=tam_def, color="blanco")
        mapping[key + "_DEF"] = str(bloque.contenido)

    elif tipo == "codigo":
        LABEL_H = 0.38 if bloque.label else 0.0
        if bloque.label:
            _texto(slide, _MX, y, 5, LABEL_H,
                   f"💻  {bloque.label}",
                   tamano=_T["label_bloque"], color="ambar", negrita=True)
        y_caja = y + LABEL_H
        h_caja = h - LABEL_H
        ancho_caja = _W - 2 * _MX
        _rect(slide, _MX, y_caja, ancho_caja, h_caja, _C["terminal_bg"])
        _rect(slide, _MX, y_caja, ancho_caja, 0.10, _C["turquesa"])
        tam_cod = calcular_tamano_codigo(str(bloque.contenido), ancho_caja - 0.30)
        _texto(slide, _MX + 0.20, y_caja + 0.15,
               ancho_caja - 0.30, h_caja - 0.25,
               "{{" + key + "_COD}}", tamano=tam_cod,
               color="codigo_color", fuente="Consolas")
        mapping[key + "_COD"] = str(bloque.contenido)

    elif tipo == "idea_clave":
        LABEL_H = 0.38
        _rect(slide, _MX, y, _W - 2 * _MX, h, _C["fondo_bloque"])
        _rect(slide, _MX, y, 0.08, h, _C["ambar"])
        _texto(slide, _MX + 0.25, y + 0.08, _W - 2 * _MX - 0.4, LABEL_H,
               f"💡  Idea clave",
               tamano=_T["titulo_bloque"], color="ambar", negrita=True)
        _texto(slide, _MX + 0.25, y + LABEL_H + 0.05,
               _W - 2 * _MX - 0.4, h - LABEL_H - 0.10,
               "{{" + key + "_IK}}", tamano=_T["cuerpo_stock"], color="blanco",
               anchor_vertical=MSO_ANCHOR.MIDDLE)
        mapping[key + "_IK"] = str(bloque.contenido)

    elif tipo == "bullets":
        LABEL_H = 0.38 if bloque.label else 0.0
        _rect(slide, _MX, y, _W - 2 * _MX, h, _C["fondo_bloque"])
        _rect(slide, _MX, y, 0.08, h, _C["turquesa"])
        if bloque.label:
            _texto(slide, _MX + 0.25, y + 0.10, _W - 2 * _MX - 0.4, LABEL_H,
                   bloque.label, tamano=_T["titulo_bloque"], color="turquesa", negrita=True)
        items = bloque.contenido if isinstance(bloque.contenido, list) else [str(bloque.contenido)]
        texto_bullets = "\n".join(f"• {item}" for item in items)
        _texto(slide, _MX + 0.25, y + LABEL_H + 0.10,
               _W - 2 * _MX - 0.4, h - LABEL_H - 0.15,
               "{{" + key + "_BUL}}", tamano=20, color="blanco")
        mapping[key + "_BUL"] = texto_bullets

    elif tipo == "tabla":
        filas = bloque.contenido if isinstance(bloque.contenido, list) else []
        insertar_tabla_pptx(slide, filas,
                            Inches(_MX), Inches(y),
                            Inches(_W - 2 * _MX), Inches(h))

    # "separador" y otros tipos desconocidos: noop


def construir_slide_icn_flexible(prs, num: str, plan: "SlidePlan"):
    """Genera un slide ICN componiendo bloques dinámicamente.

    Los bloques se distribuyen verticalmente en el área disponible (1.65"–7.20"),
    con alturas proporcionales a su costo en filas según el planificador.
    """
    if not _HELPERS_OK:
        log_warn("construir_slide_icn_flexible: helpers visuales no disponibles, "
                 "usando layout clásico como fallback")
        # fallback: usar concepto clásico si hay datos del concepto
        if plan.concepto:
            construir_concepto_clasico(prs, num, plan.concepto)
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fondo_oscuro(slide)
    _barra_superior(slide)

    # Título
    _texto(slide, _MX, 0.95, _W - 2 * _MX, 0.75,
           plan.titulo, tamano=_T["titulo_stock"], color="blanco", negrita=True)

    # Calcular alturas proporcionales
    bloques = plan.bloques
    densidad_total = sum(b.filas for b in bloques)
    n_bloques_reales = sum(1 for b in bloques if b.tipo != "separador")
    n_gaps = max(0, n_bloques_reales - 1)
    area_total = _Y_ICN_FIN - _Y_ICN_INICIO
    area_para_bloques = area_total - _GAP_BLOQUES * n_gaps
    coef = area_para_bloques / max(densidad_total, 0.1)

    if plan.justificacion:
        log_warn(f"[planificador] {plan.justificacion}")

    # Renderizar bloques
    mapping = {
        "NUMERO_CLASE": num,
        "SECCION": plan.seccion,
    }
    y = _Y_ICN_INICIO
    for idx, bloque in enumerate(bloques):
        if bloque.tipo == "separador":
            y += _GAP_BLOQUES
            continue
        h = max(bloque.filas * coef, 0.40)
        key = f"BLK{idx}"
        _renderizar_bloque(slide, bloque, y, h, key, mapping)
        y += h + _GAP_BLOQUES

    reemplazar_placeholders(slide, mapping)


# =====================================================================
# SLIDE HAZ AHORA FLEXIBLE
# =====================================================================

def construir_slide_haz_ahora_flex(prs, num: str, plan_ha: dict):
    """Genera el slide de Haz Ahora con composición dinámica.

    Diseño tipo "situaciones" (caso más común):
    - Barra superior (NUMERO_CLASE + SECCION = "Haz Ahora")
    - Título "⚡ Haz Ahora" en ámbar
    - Caja intro pequeña (la instrucción) con borde ámbar
    - Caja grande numerada con las situaciones
    - Nota de cierre al fondo en gris
    """
    if not _HELPERS_OK:
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fondo_oscuro(slide)
    _barra_superior(slide)

    mapping = {"NUMERO_CLASE": num, "SECCION": "Haz Ahora"}

    # Título
    _texto(slide, _MX, 0.62, _W - 2 * _MX, 0.52,
           "⚡ Haz Ahora",
           tamano=30, color="ambar", negrita=True)

    tipo = plan_ha.get("tipo", "libre")
    intro = plan_ha.get("intro", "")
    situaciones = plan_ha.get("situaciones", [])
    cierre = plan_ha.get("cierre", "")

    Y_START = 1.22
    Y_END = 7.22
    GAP = 0.12

    y = Y_START

    if tipo == "situaciones" and situaciones:
        # Altura dinámica del intro según cuántas líneas ocupa el texto
        n_lineas_intro = estimar_lineas(intro, 17, _W - 2 * _MX - 0.30)
        INTRO_H = max(0.70, n_lineas_intro * 0.32 + 0.28)
        CIERRE_H = 0.48

        # Caja intro
        if intro:
            _rect(slide, _MX, y, _W - 2 * _MX, INTRO_H, _C["fondo_bloque"])
            _rect(slide, _MX, y, 0.08, INTRO_H, _C["ambar"])
            _texto(slide, _MX + 0.22, y + 0.09,
                   _W - 2 * _MX - 0.30, INTRO_H - 0.14,
                   "{{HA_INTRO}}", tamano=17, color="gris_claro")
            mapping["HA_INTRO"] = intro
            y += INTRO_H + GAP

        # Espacio para cierre al fondo
        y_cierre_start = Y_END - CIERRE_H
        h_sit = (y_cierre_start - GAP) - y

        # Caja situaciones numeradas
        n_sit = len(situaciones)
        tam_sit = 20 if n_sit <= 6 else 18

        _rect(slide, _MX, y, _W - 2 * _MX, h_sit, _C["fondo_bloque"])
        _rect(slide, _MX, y, 0.08, h_sit, _C["turquesa"])
        texto_sit = "\n".join(situaciones)
        _texto(slide, _MX + 0.22, y + 0.15,
               _W - 2 * _MX - 0.30, h_sit - 0.22,
               "{{HA_SIT}}", tamano=tam_sit, color="blanco")
        mapping["HA_SIT"] = texto_sit

        # Nota de cierre
        if cierre:
            _texto(slide, _MX, y_cierre_start,
                   _W - 2 * _MX, CIERRE_H,
                   "{{HA_CIERRE}}", tamano=15, color="gris_secundario")
            mapping["HA_CIERRE"] = "→ " + cierre

    else:
        # Layout libre: todo el contenido en una sola caja grande
        h = Y_END - y
        _rect(slide, _MX, y, _W - 2 * _MX, h, _C["fondo_bloque"])
        _rect(slide, _MX, y, 0.08, h, _C["ambar"])
        texto_libre = intro
        if situaciones:
            texto_libre += ("\n\n" if intro else "") + "\n".join(situaciones)
        _texto(slide, _MX + 0.22, y + 0.15,
               _W - 2 * _MX - 0.30, h - 0.22,
               "{{HA_LIBRE}}", tamano=20, color="blanco")
        mapping["HA_LIBRE"] = texto_libre

    reemplazar_placeholders(slide, mapping)


def generar_preview_haz_ahora(spec: dict, ruta_salida: Path):
    """Genera un PPT con solo el slide de Haz Ahora."""
    if not RUTA_PLANTILLA.exists():
        raise FileNotFoundError(f"No se encontro la plantilla en {RUTA_PLANTILLA}")

    prs = Presentation(str(RUTA_PLANTILLA))
    n = spec["numero_clase"]
    num = f"{n:02d}" if isinstance(n, int) else (str(n) if n is not None else "??")

    if spec["haz_ahora"] and _PLANIFICADOR_OK and _HELPERS_OK:
        plan_ha = planificar_haz_ahora(spec["haz_ahora"])
        construir_slide_haz_ahora_flex(prs, num, plan_ha)

    eliminar_slides(prs, list(range(TOTAL_SLIDES_MODELO)))
    ruta_salida.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(ruta_salida))


# =====================================================================
# CONSTRUCCIÓN DEL PPT
# =====================================================================

def _construir_icn(prs, num: str, spec: dict):
    """Orquesta la generación de slides ICN usando el planificador."""
    conceptos = spec["icn"]["conceptos"]
    demos = spec["icn"].get("demos_apiladas", [])

    if _PLANIFICADOR_OK:
        planes = planificar_icn(conceptos, demos, num)
    else:
        # Fallback: comportamiento v4 (un slide por concepto)
        planes = _planes_fallback(conceptos, demos)

    for plan in planes:
        tipo = plan.tipo_slide

        if tipo == "icn_flexible":
            if _HELPERS_OK:
                construir_slide_icn_flexible(prs, num, plan)
            else:
                # Fallback al constructor clásico
                construir_concepto_clasico(prs, num, plan.concepto)

        elif tipo == "anatomia":
            construir_concepto_anatomia(prs, num, plan.concepto)

        elif tipo == "analogia":
            construir_concepto_analogia(prs, num, plan.concepto)

        elif tipo == "antes_despues":
            construir_concepto_antes_despues(prs, num, plan.concepto)

        elif tipo == "frase_clave":
            construir_concepto_frase_clave(prs, num, plan.concepto)

        elif tipo == "tabla_demos":
            _construir_demo_apilada(prs, num, plan.concepto)

        # Flag: slide de frase clave aparte después del concepto
        concepto = plan.concepto or {}
        if concepto.get("frase_clave_aparte") and concepto.get("idea_clave"):
            construir_concepto_frase_clave(
                prs, num,
                {"idea_clave": concepto["idea_clave"],
                 "nombre": concepto.get("nombre", ""),
                 "numero": concepto.get("numero", "")},
                solo_frase=True,
            )


def _construir_demo_apilada(prs, num: str, demo: dict):
    """Construye un slide de demostración apilada (tipo 9)."""
    sl = duplicar_slide(prs, prs.slides[M_CODIGO_RESULTADO_MULTI])
    filas = (demo.get("filas") or [])[:3]

    mapping = {
        "NUMERO_CLASE": num,
        "SECCION": "Contenido nuevo",
        "TITULO_SLIDE": demo.get("titulo", "Demostración"),
        "SUBTITULO_SLIDE": demo.get("subtitulo", ""),
    }
    for idx in range(1, 4):
        if idx <= len(filas):
            fila = filas[idx - 1]
            mapping[f"FILA_{idx}_LABEL"]     = fila.get("label", "")
            mapping[f"FILA_{idx}_CODIGO"]    = fila.get("codigo", "")
            mapping[f"FILA_{idx}_RESULTADO"] = fila.get("resultado", "")
        else:
            mapping[f"FILA_{idx}_LABEL"]     = ""
            mapping[f"FILA_{idx}_CODIGO"]    = ""
            mapping[f"FILA_{idx}_RESULTADO"] = ""

    for idx in range(1, len(filas) + 1):
        fila = filas[idx - 1]
        ajustar_fuente_codigo_en_shape(
            sl, f"FILA_{idx}_CODIGO", ANCHO_CAJA_CODIGO_MULTI,
            fila.get("codigo", ""))
        ajustar_fuente_codigo_en_shape(
            sl, f"FILA_{idx}_RESULTADO", ANCHO_CAJA_RESULTADO_MULTI,
            fila.get("resultado", ""))

    reemplazar_placeholders(sl, mapping)

    filas_vacias = [i for i in range(1, 4) if i > len(filas)]
    if filas_vacias:
        ocultar_filas_vacias_multi(sl, filas_vacias)


def _planes_fallback(conceptos: list, demos: list) -> list:
    """Genera planes con el comportamiento v4 (un slide por concepto).
    Usado si planificar_slides no está disponible.
    """
    from dataclasses import dataclass, field as _field

    @dataclass
    class _Plan:
        tipo_slide: str
        titulo: str
        concepto: dict = None
        bloques: list = _field(default_factory=list)
        seccion: str = "Contenido nuevo"
        densidad: float = 0.0
        justificacion: str = ""

    planes = []
    for c in conceptos:
        tipo_ex = (c.get("tipo") or "").lower()
        if tipo_ex in {"anatomia", "anatomía"} or c.get("partes"):
            t = "anatomia"
        elif tipo_ex in {"analogia", "analogía"} or c.get("analogia_filas"):
            t = "analogia"
        elif c.get("antes_codigo"):
            t = "antes_despues"
        elif tipo_ex in {"frase_clave"}:
            t = "frase_clave"
        else:
            t = "icn_flexible"  # usa clásico via fallback interno
        planes.append(_Plan(tipo_slide=t, titulo=c.get("nombre", ""), concepto=c))

    for demo in demos:
        marcado = dict(demo)
        marcado["_es_demo"] = True
        planes.append(_Plan(tipo_slide="tabla_demos",
                            titulo=demo.get("titulo", ""), concepto=marcado))
    return planes


def generar_preview_icn(spec: dict, ruta_salida: Path):
    """Genera un PPT con solo las slides de ICN (conceptos + demos).

    Usado para validar la sección ICN sin necesidad de ver el PPT completo.
    No incluye: bienvenida, objetivo, haz ahora, errores comunes.
    """
    if not RUTA_PLANTILLA.exists():
        raise FileNotFoundError(f"No se encontró la plantilla en {RUTA_PLANTILLA}")

    prs = Presentation(str(RUTA_PLANTILLA))
    n = spec["numero_clase"]
    num = f"{n:02d}" if isinstance(n, int) else (str(n) if n is not None else "??")

    _construir_icn(prs, num, spec)

    eliminar_slides(prs, list(range(TOTAL_SLIDES_MODELO)))
    ruta_salida.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(ruta_salida))


def generar_ppt(spec: dict, ruta_salida: Path):
    if not RUTA_PLANTILLA.exists():
        raise FileNotFoundError(f"No se encontró la plantilla en {RUTA_PLANTILLA}")

    prs = Presentation(str(RUTA_PLANTILLA))
    n = spec["numero_clase"]
    num = f"{n:02d}" if isinstance(n, int) else (str(n) if n is not None else "??")
    tema = spec["tema"] or "Clase de Python"

    # === Bienvenida ===
    sl = duplicar_slide(prs, prs.slides[M_BIENVENIDA])
    reemplazar_placeholders(sl, {
        "TITULO_CLASE": tema,
        "SUBTITULO": generar_subtitulo(spec),
        "NUMERO_CLASE": num,
        "MENSAJE_HOY": generar_mensaje_hoy(spec),
    })

    # === Objetivo ===
    sl = duplicar_slide(prs, prs.slides[M_OBJETIVO])
    reemplazar_placeholders(sl, {
        "NUMERO_CLASE": num,
        "TITULO_CLASE": tema,
        "OBJETIVO": spec["objetivo"] or "",
        "PROPOSITO": spec["proposito"] or "",
        "REGLAS": generar_reglas(),
    })

    # === Haz Ahora ===
    if spec["haz_ahora"]:
        if _PLANIFICADOR_OK and _HELPERS_OK:
            plan_ha = planificar_haz_ahora(spec["haz_ahora"])
            construir_slide_haz_ahora_flex(prs, num, plan_ha)
        else:
            # Fallback: comportamiento v4 (clon de slide modelo)
            sl = duplicar_slide(prs, prs.slides[M_STOCK_TITULO])
            ha = separar_haz_ahora(spec["haz_ahora"])
            tiene_codigo_ha = bool(ha["codigo"])
            shape_cod_ha = buscar_shape_por_key(sl, "CONTENIDO_CODIGO")
            shape_txt_ha = buscar_shape_por_key(sl, "CONTENIDO_TEXTO")
            reemplazar_placeholders(sl, {
                "NUMERO_CLASE": num,
                "SECCION": "Haz Ahora",
                "TITULO_SLIDE": "Haz Ahora",
                "SUBTITULO_SLIDE": ha["intro"],
                "CONTENIDO_TEXTO": ha["texto"],
                "CONTENIDO_CODIGO": ha["codigo"] or "",
                "NOTA_INFERIOR": "Idea clave: lo que ya sabes te ayuda a entender lo nuevo.",
            })
            if not tiene_codigo_ha and shape_cod_ha and shape_txt_ha:
                extra = shape_cod_ha.height
                eliminar_shape_de_slide(sl, shape_cod_ha)
                shape_txt_ha.height += extra

    # === ICN — planificador decide composición ===
    _construir_icn(prs, num, spec)

    # === Tabla de errores (última slide del PPT) ===
    if spec["icn"]["errores_tabla"]:
        sl = duplicar_slide(prs, prs.slides[M_TABLA])
        mapping = {
            "NUMERO_CLASE": num,
            "SECCION": "Errores típicos",
            "TITULO_SLIDE": "⚠️ Errores típicos a evitar",
            "NOTA_INFERIOR": "💡 Si te aparece uno de estos errores, vuelve a esta tabla antes de pedir ayuda.",
        }
        for i, fila in enumerate(spec["icn"]["errores_tabla"][:3], start=1):
            for j, valor in enumerate(fila[:3], start=1):
                mapping[f"FILA_{i}_COL_{j}"] = valor
        for i in range(1, 4):
            for j in range(1, 4):
                k = f"FILA_{i}_COL_{j}"
                if k not in mapping:
                    mapping[k] = ""
        reemplazar_placeholders(sl, mapping)

    # STOP: Guiada, Independiente, Ticket y Cierre se trabajan desde el Colab.

    eliminar_slides(prs, list(range(TOTAL_SLIDES_MODELO)))
    ruta_salida.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(ruta_salida))


# =====================================================================
# HELPERS DE CONTENIDO
# =====================================================================

def generar_subtitulo(spec: dict) -> str:
    contenidos_nuevos = spec["contexto"].get("Contenidos nuevos", "")
    if contenidos_nuevos:
        return f"Hoy: {contenidos_nuevos}"
    return ""


def generar_mensaje_hoy(spec: dict) -> str:
    """Mensaje motivador para la portada.

    Combina el tema de la clase con un mensaje inspirador relacionado.
    Si el propósito tiene una analogía o ejemplo concreto (Steam, juegos,
    apps, etc.), lo usa. Si no, genera uno genérico pero específico al tema.
    """
    tema = spec.get("tema") or "programar"
    proposito = spec.get("proposito") or ""

    # Intentar extraer la frase con la analogía/ejemplo concreto del propósito
    # (típicamente después de "—" o que mencione una app/juego/situación real)
    if proposito:
        # Buscar oración que tenga una analogía concreta (después de — o que mencione "como")
        match_analogia = re.search(
            r"(?:—|—|como )\s*([^.—]+(?:Steam|juego|app|aplicación|YouTube|Spotify|Instagram|TikTok|servicio|sitio|red social|programa)[^.]*)",
            proposito, re.IGNORECASE
        )
        if match_analogia:
            analogia = match_analogia.group(1).strip().rstrip(",;:")
            if len(analogia) < 130:
                return f"Hoy darás un salto: pasarás de ver código a controlarlo {analogia} 🚀"

        # Si no hay analogía clara, usar la primera oración del propósito pero con marco motivador
        primera = proposito.split(".")[0].strip()
        if len(primera) > 0 and len(primera) < 130:
            # Si la primera oración suena descriptiva, le agregamos un marco motivador
            return f"{primera}. Hoy desbloqueas una herramienta nueva ⚡"

    # Fallback: mensaje genérico pero específico al tema
    return f"Hoy aprenderás {tema} — un paso más para que tu código haga cosas reales 🚀"


def generar_reglas() -> str:
    return ("📓 Todo comando nuevo va al cuaderno.\n"
            "🚫 NO hagan copiar-pegar de los comandos.\n"
            "✨ En toda evaluación podrán usar sus apuntes.")


# =====================================================================
# MAIN
# =====================================================================

def main():
    args = sys.argv[1:]
    modo_preview_icn    = "--preview-icn"      in args
    modo_preview_ha     = "--preview-hazahora" in args
    args_limpios = [a for a in args if not a.startswith("--")]

    if len(args_limpios) != 2:
        print("Uso: python crear_ppt.py <ruta_spec.md> <ruta_salida.pptx> [--preview-icn | --preview-hazahora]")
        print()
        print("  --preview-icn        genera solo las slides de ICN")
        print("  --preview-hazahora   genera solo el slide de Haz Ahora")
        sys.exit(1)

    ruta_spec   = Path(args_limpios[0])
    ruta_salida = Path(args_limpios[1])

    if not ruta_spec.exists():
        print(f"ERROR: no existe {ruta_spec}")
        sys.exit(1)
    if not RUTA_PLANTILLA.exists():
        print(f"ERROR: falta plantilla en {RUTA_PLANTILLA}")
        sys.exit(1)

    print(f"[spec] Leyendo: {ruta_spec}")
    spec = parsear_spec(ruta_spec)
    print(f"[ok]   Clase {spec['numero_clase']} - {spec['tema']}")
    print(f"       {len(spec['icn']['conceptos'])} conceptos ICN, "
          f"{len(spec['icn'].get('demos_apiladas', []))} demos, "
          f"{len(spec['icn']['errores_tabla'])} errores tipicos")

    if modo_preview_icn:
        print("[preview] Generando solo slides de ICN...")
        generar_preview_icn(spec, ruta_salida)
        print(f"[ok] Preview ICN guardado en: {ruta_salida}")
    elif modo_preview_ha:
        print("[preview] Generando solo slide de Haz Ahora...")
        generar_preview_haz_ahora(spec, ruta_salida)
        print(f"[ok] Preview Haz Ahora guardado en: {ruta_salida}")
    else:
        print("[ppt] Generando presentacion completa (hasta errores comunes)...")
        generar_ppt(spec, ruta_salida)
        print(f"[ok] PPT guardado en: {ruta_salida}")


if __name__ == "__main__":
    main()
