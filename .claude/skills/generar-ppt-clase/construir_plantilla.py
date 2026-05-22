#!/usr/bin/env python3
"""
construir_plantilla.py — Construye plantilla_marca.pptx desde cero. (v4)

Cambios v4:
- Nuevo slide modelo (índice 9): "Código + Resultado apilado" (multi-fila),
  el layout estrella para contenido denso/compacto. Hasta 3 filas, cada una
  con etiqueta + código (izq) + resultado (der).
- Helper calcular_tamano_codigo(): base de la regla mixta de fuente para
  código (el ajuste real sobre contenido se aplica en crear_ppt.py).

Cambios v3:
- Fondo oscuro garantizado con rectángulo full-slide (más fiable que slide.background)
- Cajas Objetivo/Propósito más altas, Reglas más baja
- En Stock Título y Haz Ahora, dos zonas: texto narrativo arriba + bloque código terminal opcional abajo
- En slides de concepto, Definición e Idea Clave más altas
- Tabla con texto centrado verticalmente; encabezados centrados horizontalmente, celdas a la izquierda
- Subtítulos en blanco más grandes
- Emojis identificadores por sección
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

# =====================================================================
# PALETA DE MARCA
# =====================================================================

C = {
    "fondo_oscuro": RGBColor(0x1A, 0x1A, 0x2E),
    "fondo_bloque": RGBColor(0x16, 0x21, 0x3E),
    "turquesa": RGBColor(0x0F, 0x9B, 0x8E),
    "menta": RGBColor(0x02, 0xC3, 0x9A),
    "ambar": RGBColor(0xF5, 0xA6, 0x23),
    "blanco": RGBColor(0xFF, 0xFF, 0xFF),
    "gris_claro": RGBColor(0xCB, 0xD5, 0xE1),
    "gris_secundario": RGBColor(0x94, 0xA3, 0xB8),
    "terminal_bg": RGBColor(0x0D, 0x11, 0x17),
    "codigo_color": RGBColor(0x4A, 0xDF, 0xCB),
    "tabla_header": RGBColor(0x0F, 0x16, 0x2A),
    "tabla_fila": RGBColor(0x1F, 0x29, 0x44),
}

T = {
    "prompt_grande": 76,
    "titulo_bienvenida": 40,
    "subtitulo_bienvenida": 20,
    "metadatos": 14,
    "header_objetivo": 16,
    "titulo_bloque": 18,
    "cuerpo_bloque": 22,
    "titulo_stock": 32,
    "subtitulo_stock": 22,
    "cuerpo_stock": 22,
    "label_bloque": 18,
    "codigo": 20,
    "barra_navegacion": 12,
    "nota_inferior": 20,
    "tabla_header": 18,
    "tabla_celda": 16,
}

# Emojis por sección
E = {
    "haz_ahora": "⚡",
    "icn": "📘",
    "errores": "⚠️",
    "guiada": "🛠️",
    "independiente": "💻",
    "ticket": "🎫",
    "cierre": "🧠",
    "idea": "💡",
}

ANCHO_SLIDE = 13.333
ALTO_SLIDE = 7.5
MARGEN_X = 0.6
BARRA_ALTO = 0.55


# =====================================================================
# HELPERS BASE
# =====================================================================

def aplicar_fondo_oscuro(slide):
    """Aplica fondo oscuro garantizado mediante un rectángulo full-slide.

    Es más fiable que slide.background.fill porque algunos visores/PowerPoint
    no respetan el background del slide cuando se duplica programáticamente.
    """
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0,
        Inches(ANCHO_SLIDE), Inches(ALTO_SLIDE)
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = C["fondo_oscuro"]
    rect.line.fill.background()
    rect.shadow.inherit = False
    # Asegurar que va al fondo (z-order más bajo)
    spTree = rect._element.getparent()
    spTree.remove(rect._element)
    # Insertar justo después de nvGrpSpPr (el primer elemento del spTree)
    spTree.insert(2, rect._element)


def agregar_rectangulo(slide, x, y, w, h, color_relleno, color_borde=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color_relleno
    if color_borde is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = color_borde
    shape.shadow.inherit = False
    return shape


def agregar_texto(slide, x, y, w, h, texto, tamano=20, color="blanco",
                  negrita=False, fuente="Calibri", alineacion=PP_ALIGN.LEFT,
                  anchor_vertical=MSO_ANCHOR.TOP, auto_size=False):
    cuadro = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = cuadro.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor_vertical
    if auto_size:
        from pptx.enum.text import MSO_AUTO_SIZE
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)

    p = tf.paragraphs[0]
    p.alignment = alineacion
    run = p.add_run()
    run.text = texto
    run.font.name = fuente
    run.font.size = Pt(tamano)
    run.font.bold = negrita
    if isinstance(color, str):
        run.font.color.rgb = C[color]
    else:
        run.font.color.rgb = color
    return cuadro


def calcular_tamano_codigo(texto, ancho_caja_in, tam_max=None, tam_min=16):
    """Regla mixta de fuente para bloques de código (Consolas).

    Mide la línea más larga del código y elige el mayor tamaño de fuente
    (entre tam_min y tam_max) con el que esa línea aún cabe en el ancho de
    la caja. Resuelve el desborde de raíz en vez de delegar al auto_size
    flojo de PowerPoint.

    - Consolas es monoespaciada: cada glifo ocupa ~0.60 * tamaño_pt de ancho.
    - ancho_caja_in es el ancho útil en pulgadas (ya descontado el padding).
    - Escalones probados: tam_max, 18, tam_min (16). Devuelve el primero que cabe.
    """
    if tam_max is None:
        tam_max = T["codigo"]
    lineas = (texto or "").split("\n")
    max_chars = max((len(l) for l in lineas), default=0)
    if max_chars == 0:
        return tam_max
    ancho_util_pt = (ancho_caja_in - 0.30) * 72.0  # padding ~0.15" por lado
    # Escalones de mayor a menor
    escalones = sorted({tam_max, 18, tam_min}, reverse=True)
    for tam in escalones:
        ancho_glifo = 0.60 * tam
        if max_chars * ancho_glifo <= ancho_util_pt:
            return tam
    return tam_min


def agregar_franja_lateral(slide):
    agregar_rectangulo(slide, 0, 0, 0.12, ALTO_SLIDE, C["turquesa"])


def agregar_barra_superior(slide):
    agregar_rectangulo(slide, 0, 0, ANCHO_SLIDE, BARRA_ALTO, C["tabla_header"])
    agregar_texto(slide, MARGEN_X - 0.15, 0.13, 6, 0.35,
                  "🐍 Python · Clase {{NUMERO_CLASE}}",
                  tamano=T["barra_navegacion"], color="gris_secundario",
                  fuente="Consolas", anchor_vertical=MSO_ANCHOR.MIDDLE)
    agregar_texto(slide, ANCHO_SLIDE - 4, 0.13, 3.5, 0.35, "{{SECCION}}",
                  tamano=T["barra_navegacion"], color="turquesa",
                  fuente="Consolas", negrita=True,
                  alineacion=PP_ALIGN.RIGHT, anchor_vertical=MSO_ANCHOR.MIDDLE)


def agregar_titulo_slide(slide, texto_placeholder="{{TITULO_SLIDE}}", y=0.95):
    agregar_texto(slide, MARGEN_X, y, ANCHO_SLIDE - 2 * MARGEN_X, 0.75,
                  texto_placeholder, tamano=T["titulo_stock"],
                  color="blanco", negrita=True, fuente="Calibri")


def agregar_bloque_codigo(slide, x, y, w, h, placeholder, label=None,
                          label_color="ambar"):
    if label:
        agregar_texto(slide, x, y, w, 0.40, label,
                      tamano=T["label_bloque"], color=label_color,
                      negrita=True, fuente="Calibri")
        y_bloque = y + 0.45
        h_bloque = h - 0.45
    else:
        y_bloque = y
        h_bloque = h
    agregar_rectangulo(slide, x, y_bloque, w, h_bloque, C["terminal_bg"])
    agregar_rectangulo(slide, x, y_bloque, w, 0.10, C["turquesa"])
    agregar_texto(slide, x + 0.15, y_bloque + 0.18, w - 0.3, h_bloque - 0.25,
                  placeholder, tamano=T["codigo"], color="codigo_color",
                  fuente="Consolas")


def agregar_bloque_terminal_resultado(slide, x, y, w, h, placeholder, label=None):
    if label:
        agregar_texto(slide, x, y, w, 0.40, label,
                      tamano=T["label_bloque"], color="menta",
                      negrita=True, fuente="Calibri")
        y_bloque = y + 0.45
        h_bloque = h - 0.45
    else:
        y_bloque = y
        h_bloque = h
    agregar_rectangulo(slide, x, y_bloque, w, h_bloque, C["fondo_bloque"])
    agregar_rectangulo(slide, x, y_bloque, w, 0.10, C["menta"])
    agregar_texto(slide, x + 0.15, y_bloque + 0.18, w - 0.3, h_bloque - 0.25,
                  placeholder, tamano=T["codigo"], color="blanco",
                  fuente="Consolas")


def agregar_nota_inferior(slide, y=6.5):
    agregar_rectangulo(slide, MARGEN_X, y, ANCHO_SLIDE - 2 * MARGEN_X, 0.7,
                       C["fondo_bloque"])
    agregar_rectangulo(slide, MARGEN_X, y, 0.08, 0.7, C["turquesa"])
    agregar_texto(slide, MARGEN_X + 0.2, y + 0.05,
                  ANCHO_SLIDE - 2 * MARGEN_X - 0.3, 0.6, "{{NOTA_INFERIOR}}",
                  tamano=T["nota_inferior"], color="blanco",
                  fuente="Calibri", anchor_vertical=MSO_ANCHOR.MIDDLE)


# =====================================================================
# SLIDES MODELO
# =====================================================================

def construir_slide_bienvenida(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    aplicar_fondo_oscuro(slide)
    agregar_franja_lateral(slide)

    agregar_texto(slide, 0.75, 0.65, 4, 1.0, ">>>",
                  tamano=T["prompt_grande"], color="turquesa",
                  fuente="Consolas", negrita=True)
    agregar_texto(slide, 0.75, 1.70, ANCHO_SLIDE - 1.5, 1.0,
                  "{{TITULO_CLASE}}", tamano=T["titulo_bienvenida"],
                  color="blanco", negrita=True)
    agregar_texto(slide, 0.75, 2.72, ANCHO_SLIDE - 1.5, 0.5,
                  "{{SUBTITULO}}", tamano=T["subtitulo_bienvenida"],
                  color="gris_secundario")
    agregar_texto(slide, 0.75, 3.32, ANCHO_SLIDE - 1.5, 0.4,
                  "Clase {{NUMERO_CLASE}}", tamano=T["metadatos"],
                  color="turquesa", negrita=True)
    agregar_rectangulo(slide, 0.75, 4.10, ANCHO_SLIDE - 1.5, 2.10, C["fondo_bloque"])
    agregar_rectangulo(slide, 0.75, 4.10, 0.08, 2.10, C["ambar"])
    agregar_texto(slide, 1.0, 4.18, ANCHO_SLIDE - 2, 0.4, "Mensaje de hoy",
                  tamano=T["titulo_bloque"], color="ambar", negrita=True)
    agregar_texto(slide, 1.0, 4.65, ANCHO_SLIDE - 2, 1.45,
                  "{{MENSAJE_HOY}}", tamano=T["cuerpo_bloque"], color="blanco")


def construir_slide_objetivo(prs):
    """Objetivo/Propósito MÁS GRANDES, Reglas más chica."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    aplicar_fondo_oscuro(slide)
    agregar_franja_lateral(slide)

    agregar_texto(slide, 0.75, 0.30, ANCHO_SLIDE - 1.5, 0.45,
                  ">>> Clase {{NUMERO_CLASE}} - {{TITULO_CLASE}}",
                  tamano=T["header_objetivo"], color="turquesa",
                  fuente="Consolas")

    # Objetivo y Propósito: altura aumentada de 2.10 a 3.10
    OBJ_Y = 1.00
    OBJ_H = 3.10

    agregar_rectangulo(slide, 0.75, OBJ_Y, 5.85, OBJ_H, C["fondo_bloque"])
    agregar_rectangulo(slide, 0.75, OBJ_Y, 0.08, OBJ_H, C["turquesa"])
    agregar_texto(slide, 1.0, OBJ_Y + 0.10, 5.5, 0.40,
                  f"{E['icn']}  Objetivo",
                  tamano=T["titulo_bloque"], color="turquesa", negrita=True)
    agregar_texto(slide, 1.0, OBJ_Y + 0.55, 5.5, OBJ_H - 0.65,
                  "{{OBJETIVO}}", tamano=T["cuerpo_bloque"], color="blanco",
                  auto_size=True)

    agregar_rectangulo(slide, 6.85, OBJ_Y, 5.85, OBJ_H, C["fondo_bloque"])
    agregar_rectangulo(slide, 6.85, OBJ_Y, 0.08, OBJ_H, C["ambar"])
    agregar_texto(slide, 7.1, OBJ_Y + 0.10, 5.5, 0.40,
                  f"{E['idea']}  Propósito",
                  tamano=T["titulo_bloque"], color="ambar", negrita=True)
    agregar_texto(slide, 7.1, OBJ_Y + 0.55, 5.5, OBJ_H - 0.65,
                  "{{PROPOSITO}}", tamano=T["cuerpo_bloque"], color="blanco",
                  auto_size=True)

    # Reglas: altura reducida de 3.50 a 2.60
    REG_Y = OBJ_Y + OBJ_H + 0.20
    REG_H = ALTO_SLIDE - REG_Y - 0.30  # margen inferior

    agregar_rectangulo(slide, 0.75, REG_Y, ANCHO_SLIDE - 1.5, REG_H,
                       C["fondo_bloque"])
    agregar_rectangulo(slide, 0.75, REG_Y, 0.08, REG_H, C["menta"])
    agregar_texto(slide, 1.0, REG_Y + 0.10, ANCHO_SLIDE - 2, 0.45,
                  "☝ Reglas de ahora en adelante",
                  tamano=T["titulo_bloque"], color="menta", negrita=True)
    agregar_texto(slide, 1.0, REG_Y + 0.60, ANCHO_SLIDE - 2, REG_H - 0.70,
                  "{{REGLAS}}", tamano=T["cuerpo_bloque"], color="blanco",
                  auto_size=True)


def construir_slide_stock_titulo(prs):
    """Stock título con dos zonas: texto narrativo arriba + bloque código abajo (opcional).

    El placeholder CONTENIDO_TEXTO recibe el texto sin código.
    El placeholder CONTENIDO_CODIGO recibe el código (en terminal verdosa) o queda vacío.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    aplicar_fondo_oscuro(slide)
    agregar_barra_superior(slide)
    agregar_titulo_slide(slide)

    # Subtítulo en blanco más grande
    agregar_texto(slide, MARGEN_X, 1.65, ANCHO_SLIDE - 2 * MARGEN_X, 0.55,
                  "{{SUBTITULO_SLIDE}}", tamano=T["subtitulo_stock"],
                  color="blanco")

    # Zona 1: texto narrativo (bloque oscuro suave)
    agregar_rectangulo(slide, MARGEN_X, 2.30, ANCHO_SLIDE - 2 * MARGEN_X, 1.55,
                       C["fondo_bloque"])
    agregar_rectangulo(slide, MARGEN_X, 2.30, 0.08, 1.55, C["turquesa"])
    agregar_texto(slide, MARGEN_X + 0.25, 2.40,
                  ANCHO_SLIDE - 2 * MARGEN_X - 0.4, 1.35,
                  "{{CONTENIDO_TEXTO}}", tamano=T["cuerpo_stock"],
                  color="blanco", auto_size=True)

    # Zona 2: bloque de código terminal opcional
    agregar_texto(slide, MARGEN_X, 4.00, 4, 0.35,
                  f"{E['independiente']}  Código",
                  tamano=T["label_bloque"], color="ambar", negrita=True)
    agregar_rectangulo(slide, MARGEN_X, 4.40, ANCHO_SLIDE - 2 * MARGEN_X, 1.90,
                       C["terminal_bg"])
    agregar_rectangulo(slide, MARGEN_X, 4.40, ANCHO_SLIDE - 2 * MARGEN_X, 0.10,
                       C["turquesa"])
    agregar_texto(slide, MARGEN_X + 0.2, 4.55,
                  ANCHO_SLIDE - 2 * MARGEN_X - 0.3, 1.70,
                  "{{CONTENIDO_CODIGO}}", tamano=T["codigo"],
                  color="codigo_color", fuente="Consolas")

    agregar_nota_inferior(slide)


def construir_slide_concepto(prs):
    """Concepto del ICN: Definición + Ejemplo + Idea clave.

    Cajas Definición e Idea Clave más grandes; Ejemplo compacto.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    aplicar_fondo_oscuro(slide)
    agregar_barra_superior(slide)

    agregar_texto(slide, MARGEN_X, 0.85, ANCHO_SLIDE - 2 * MARGEN_X, 0.70,
                  "{{TITULO_CONCEPTO}}", tamano=T["titulo_stock"],
                  color="blanco", negrita=True)

    # Definición: altura 1.80 (era 1.30)
    DEF_Y = 1.65
    DEF_H = 1.80
    agregar_rectangulo(slide, MARGEN_X, DEF_Y,
                       ANCHO_SLIDE - 2 * MARGEN_X, DEF_H, C["fondo_bloque"])
    agregar_rectangulo(slide, MARGEN_X, DEF_Y, 0.08, DEF_H, C["turquesa"])
    agregar_texto(slide, MARGEN_X + 0.25, DEF_Y + 0.10,
                  ANCHO_SLIDE - 2 * MARGEN_X - 0.4, 0.40,
                  f"{E['icn']}  Definición",
                  tamano=T["titulo_bloque"], color="turquesa", negrita=True)
    agregar_texto(slide, MARGEN_X + 0.25, DEF_Y + 0.55,
                  ANCHO_SLIDE - 2 * MARGEN_X - 0.4, DEF_H - 0.65,
                  "{{DEFINICION}}", tamano=T["cuerpo_stock"],
                  color="blanco", auto_size=True)

    # Ejemplo: altura 1.55 (era 2.10), suficiente para 3-5 líneas de código
    EJ_Y = DEF_Y + DEF_H + 0.20
    EJ_H = 1.55
    agregar_texto(slide, MARGEN_X, EJ_Y, 4, 0.35,
                  f"{E['independiente']}  Ejemplo",
                  tamano=T["label_bloque"], color="ambar", negrita=True)
    EJ_Y2 = EJ_Y + 0.40
    EJ_H2 = EJ_H - 0.40
    agregar_rectangulo(slide, MARGEN_X, EJ_Y2,
                       ANCHO_SLIDE - 2 * MARGEN_X, EJ_H2, C["terminal_bg"])
    agregar_rectangulo(slide, MARGEN_X, EJ_Y2,
                       ANCHO_SLIDE - 2 * MARGEN_X, 0.08, C["turquesa"])
    agregar_texto(slide, MARGEN_X + 0.2, EJ_Y2 + 0.15,
                  ANCHO_SLIDE - 2 * MARGEN_X - 0.3, EJ_H2 - 0.25,
                  "{{EJEMPLO_CODIGO}}", tamano=T["codigo"],
                  color="codigo_color", fuente="Consolas")

    # Idea clave: altura 1.65 (era 1.20)
    IDEA_Y = EJ_Y + EJ_H + 0.20
    IDEA_H = ALTO_SLIDE - IDEA_Y - 0.25
    agregar_rectangulo(slide, MARGEN_X, IDEA_Y,
                       ANCHO_SLIDE - 2 * MARGEN_X, IDEA_H, C["fondo_bloque"])
    agregar_rectangulo(slide, MARGEN_X, IDEA_Y, 0.08, IDEA_H, C["ambar"])
    agregar_texto(slide, MARGEN_X + 0.25, IDEA_Y + 0.10,
                  ANCHO_SLIDE - 2 * MARGEN_X - 0.4, 0.40,
                  f"{E['idea']}  Idea clave",
                  tamano=T["titulo_bloque"], color="ambar", negrita=True)
    agregar_texto(slide, MARGEN_X + 0.25, IDEA_Y + 0.55,
                  ANCHO_SLIDE - 2 * MARGEN_X - 0.4, IDEA_H - 0.65,
                  "{{IDEA_CLAVE}}", tamano=T["cuerpo_stock"],
                  color="blanco", auto_size=True)


def construir_slide_codigo_resultado(prs):
    """Práctica guiada: subtítulo blanco grande arriba, bloques código/resultado abajo."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    aplicar_fondo_oscuro(slide)
    agregar_barra_superior(slide)
    agregar_titulo_slide(slide)

    # Subtítulo en BLANCO más grande, alto aumentado para 2-3 líneas
    agregar_texto(slide, MARGEN_X, 1.70, ANCHO_SLIDE - 2 * MARGEN_X, 1.20,
                  "{{SUBTITULO_SLIDE}}", tamano=T["subtitulo_stock"],
                  color="blanco", auto_size=True)

    # Bloques bajan: empieza en y=3.10 (antes 2.25)
    agregar_bloque_codigo(slide, MARGEN_X, 3.10, 6.4, 2.95,
                          "{{CODIGO}}",
                          label=f"{E['independiente']}  Código",
                          label_color="ambar")
    agregar_bloque_terminal_resultado(slide, MARGEN_X + 6.55, 3.10, 6.0, 2.95,
                                       "{{RESULTADO}}",
                                       label=f"{E['idea']}  Resultado")

    agregar_nota_inferior(slide)


def construir_slide_codigo_resultado_multi(prs):
    """Código + Resultado APILADO: hasta 3 filas (etiqueta + código izq + resultado der).

    Layout estrella para contenido denso/compacto. Cada fila tiene:
      - una etiqueta corta opcional (ej. "sep", "end") en negrita;
      - una caja de código a la izquierda (~6.8" de ancho);
      - una caja de resultado a la derecha (~4.8" de ancho).

    Las filas que el spec no rellene quedan con placeholders que crear_ppt.py
    vacía y oculta. La regla mixta de fuente se aplica en crear_ppt.py sobre
    el contenido real.

    Placeholders por fila i (1..3):
      {{FILA_i_LABEL}}  {{FILA_i_CODIGO}}  {{FILA_i_RESULTADO}}
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    aplicar_fondo_oscuro(slide)
    agregar_barra_superior(slide)
    agregar_titulo_slide(slide)

    # Subtítulo compacto (una línea)
    agregar_texto(slide, MARGEN_X, 1.62, ANCHO_SLIDE - 2 * MARGEN_X, 0.50,
                  "{{SUBTITULO_SLIDE}}", tamano=T["subtitulo_stock"],
                  color="blanco")

    # Geometría de las filas
    COD_X = MARGEN_X
    COD_W = 6.85
    RES_X = MARGEN_X + COD_W + 0.30
    RES_W = ANCHO_SLIDE - MARGEN_X - RES_X
    FILA_Y0 = 2.25          # inicio de la primera fila
    FILA_H = 1.42           # alto de caja de cada fila
    LABEL_H = 0.34          # alto de la etiqueta
    FILA_PASO = 1.62        # separación vertical entre filas (caja + etiqueta + gap)

    for i in range(1, 4):
        y_label = FILA_Y0 + (i - 1) * FILA_PASO
        y_caja = y_label + LABEL_H

        # Etiqueta de la fila (negrita, ámbar)
        agregar_texto(slide, COD_X, y_label, ANCHO_SLIDE - 2 * MARGEN_X, LABEL_H,
                      "{{FILA_" + str(i) + "_LABEL}}",
                      tamano=T["label_bloque"], color="ambar", negrita=True)

        # Caja de código (terminal verdosa, borde superior turquesa)
        agregar_rectangulo(slide, COD_X, y_caja, COD_W, FILA_H, C["terminal_bg"])
        agregar_rectangulo(slide, COD_X, y_caja, COD_W, 0.08, C["turquesa"])
        agregar_texto(slide, COD_X + 0.15, y_caja + 0.14, COD_W - 0.3, FILA_H - 0.22,
                      "{{FILA_" + str(i) + "_CODIGO}}",
                      tamano=T["codigo"], color="codigo_color", fuente="Consolas")

        # Caja de resultado (fondo bloque, borde superior menta)
        agregar_rectangulo(slide, RES_X, y_caja, RES_W, FILA_H, C["fondo_bloque"])
        agregar_rectangulo(slide, RES_X, y_caja, RES_W, 0.08, C["menta"])
        agregar_texto(slide, RES_X + 0.15, y_caja + 0.14, RES_W - 0.3, FILA_H - 0.22,
                      "{{FILA_" + str(i) + "_RESULTADO}}",
                      tamano=T["codigo"], color="blanco", fuente="Consolas")


def construir_slide_tabla(prs):
    """Tabla con texto centrado verticalmente, encabezados centrados horizontalmente,
    celdas alineadas a la izquierda."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    aplicar_fondo_oscuro(slide)
    agregar_barra_superior(slide)
    agregar_titulo_slide(slide)

    rows, cols = 4, 3
    x, y, w, h = MARGEN_X, 1.75, ANCHO_SLIDE - 2 * MARGEN_X, 4.50
    tabla_shape = slide.shapes.add_table(
        rows, cols, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    tabla = tabla_shape.table

    tabla.columns[0].width = Inches(3.5)
    tabla.columns[1].width = Inches(4.5)
    tabla.columns[2].width = Inches(ANCHO_SLIDE - 2 * MARGEN_X - 3.5 - 4.5)

    # Header: centrado horizontal y vertical
    encabezados = ["Error", "Qué ocurre", "Cómo corregirlo"]
    for j, enc in enumerate(encabezados):
        cell = tabla.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = C["tabla_header"]
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame
        tf.clear()
        tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.15)
        tf.margin_top = Inches(0.1)
        tf.margin_bottom = Inches(0.1)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = enc
        r.font.name = "Calibri"
        r.font.size = Pt(T["tabla_header"])
        r.font.bold = True
        r.font.color.rgb = C["blanco"]

    # Filas de datos: centrado vertical, alineado a la izquierda
    for i in range(1, rows):
        for j in range(cols):
            cell = tabla.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = C["tabla_fila"]
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame
            tf.clear()
            tf.margin_left = Inches(0.15)
            tf.margin_right = Inches(0.15)
            tf.margin_top = Inches(0.1)
            tf.margin_bottom = Inches(0.1)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = "{{FILA_" + str(i) + "_COL_" + str(j + 1) + "}}"
            r.font.name = "Calibri"
            r.font.size = Pt(T["tabla_celda"])
            r.font.color.rgb = C["blanco"]

    agregar_nota_inferior(slide)


def construir_slide_ejercicio(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    aplicar_fondo_oscuro(slide)
    agregar_barra_superior(slide)
    agregar_titulo_slide(slide)

    agregar_texto(slide, MARGEN_X, 1.70, ANCHO_SLIDE - 2 * MARGEN_X, 0.50,
                  "Resuelve este ejercicio de forma individual.",
                  tamano=T["subtitulo_stock"], color="blanco")

    agregar_rectangulo(slide, MARGEN_X, 2.30, 7.4, 3.40, C["fondo_bloque"])
    agregar_rectangulo(slide, MARGEN_X, 2.30, 0.08, 3.40, C["turquesa"])
    agregar_texto(slide, MARGEN_X + 0.25, 2.40, 6.95, 0.40,
                  f"{E['icn']}  Enunciado", tamano=T["label_bloque"],
                  color="turquesa", negrita=True)
    agregar_texto(slide, MARGEN_X + 0.25, 2.85, 6.95, 2.75,
                  "{{ENUNCIADO}}", tamano=T["cuerpo_stock"], color="blanco",
                  auto_size=True)

    agregar_bloque_terminal_resultado(slide, MARGEN_X + 7.55, 2.30, 4.85, 3.40,
                                       "{{RESULTADO_ESPERADO}}",
                                       label=f"{E['idea']}  Resultado esperado")

    agregar_nota_inferior(slide)


def construir_slide_ticket(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    aplicar_fondo_oscuro(slide)
    agregar_barra_superior(slide)

    agregar_texto(slide, MARGEN_X, 0.85, ANCHO_SLIDE - 2 * MARGEN_X, 0.75,
                  f"{E['ticket']}  Ticket de salida",
                  tamano=T["titulo_stock"], color="ambar", negrita=True)

    agregar_texto(slide, MARGEN_X, 1.60, ANCHO_SLIDE - 2 * MARGEN_X, 0.50,
                  "Tu evidencia individual de la clase de hoy.",
                  tamano=T["subtitulo_stock"], color="blanco")

    agregar_rectangulo(slide, MARGEN_X, 2.30, ANCHO_SLIDE - 2 * MARGEN_X, 2.60,
                       C["fondo_bloque"])
    agregar_rectangulo(slide, MARGEN_X, 2.30, 0.08, 2.60, C["ambar"])
    agregar_texto(slide, MARGEN_X + 0.25, 2.42,
                  ANCHO_SLIDE - 2 * MARGEN_X - 0.4, 0.40,
                  f"{E['independiente']}  Tu tarea",
                  tamano=T["titulo_bloque"], color="ambar", negrita=True)
    agregar_texto(slide, MARGEN_X + 0.25, 2.85,
                  ANCHO_SLIDE - 2 * MARGEN_X - 0.4, 2.00,
                  "{{TAREA}}", tamano=T["cuerpo_stock"], color="blanco",
                  auto_size=True)

    agregar_rectangulo(slide, MARGEN_X, 5.10, ANCHO_SLIDE - 2 * MARGEN_X, 1.65,
                       C["fondo_bloque"])
    agregar_rectangulo(slide, MARGEN_X, 5.10, 0.08, 1.65, C["menta"])
    agregar_texto(slide, MARGEN_X + 0.25, 5.22,
                  ANCHO_SLIDE - 2 * MARGEN_X - 0.4, 0.40,
                  f"{E['idea']}  Cómo entregar",
                  tamano=T["titulo_bloque"], color="menta", negrita=True)
    agregar_texto(slide, MARGEN_X + 0.25, 5.65,
                  ANCHO_SLIDE - 2 * MARGEN_X - 0.4, 1.05,
                  "{{ENTREGA}}", tamano=T["cuerpo_stock"], color="blanco",
                  auto_size=True)


def construir_slide_cierre(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    aplicar_fondo_oscuro(slide)
    agregar_barra_superior(slide)
    agregar_titulo_slide(slide, texto_placeholder=f"{E['cierre']}  Antes de irte, piensa…")

    agregar_texto(slide, MARGEN_X, 1.70, ANCHO_SLIDE - 2 * MARGEN_X, 0.50,
                  "Tres preguntas para llevarte de esta clase:",
                  tamano=T["subtitulo_stock"], color="blanco")

    y_actual = 2.40
    for i in range(1, 4):
        agregar_rectangulo(slide, MARGEN_X, y_actual,
                           ANCHO_SLIDE - 2 * MARGEN_X, 1.30, C["fondo_bloque"])
        agregar_rectangulo(slide, MARGEN_X, y_actual, 0.08, 1.30, C["turquesa"])
        agregar_texto(slide, MARGEN_X + 0.25, y_actual + 0.12,
                      0.6, 0.6, f"{i}.", tamano=T["cuerpo_stock"],
                      color="turquesa", negrita=True)
        agregar_texto(slide, MARGEN_X + 0.85, y_actual + 0.15,
                      ANCHO_SLIDE - 2 * MARGEN_X - 1.0, 1.00,
                      "{{PREGUNTA_" + str(i) + "}}",
                      tamano=T["cuerpo_stock"], color="blanco", auto_size=True)
        y_actual += 1.45


# =====================================================================
# MAIN
# =====================================================================

def main():
    prs = Presentation()
    prs.slide_width = Inches(ANCHO_SLIDE)
    prs.slide_height = Inches(ALTO_SLIDE)
    if len(prs.slides) > 0:
        for sld_id in list(prs.slides._sldIdLst):
            prs.slides._sldIdLst.remove(sld_id)

    construir_slide_bienvenida(prs)            # 0
    construir_slide_objetivo(prs)               # 1
    construir_slide_stock_titulo(prs)           # 2
    construir_slide_concepto(prs)               # 3
    construir_slide_codigo_resultado(prs)       # 4
    construir_slide_tabla(prs)                  # 5
    construir_slide_ejercicio(prs)              # 6
    construir_slide_ticket(prs)                 # 7
    construir_slide_cierre(prs)                 # 8
    construir_slide_codigo_resultado_multi(prs) # 9 — apilado (layout estrella)

    ruta_salida = Path(__file__).parent / "plantilla_marca.pptx"
    prs.save(str(ruta_salida))
    print(f"✅ Plantilla guardada en {ruta_salida}")
    print(f"   {len(prs.slides)} slides modelo")


if __name__ == "__main__":
    main()
